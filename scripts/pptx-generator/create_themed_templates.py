#!/usr/bin/env python3
"""
Generate themed PPTX templates by modifying theme XML.

This script creates properly styled template PPTX files by:
1. Creating a base presentation with python-pptx (for standard layouts)
2. Extracting the PPTX (ZIP format)
3. Modifying ppt/theme/theme1.xml with custom colors using lxml
4. Repackaging as PPTX

Usage:
    python create_themed_templates.py [--template NAME]

    Without arguments, regenerates all templates from registry.json.
    With --template, regenerates only the specified template.
"""

import argparse
import json
import logging
import os
import shutil
import tempfile
import zipfile
from pathlib import Path
from typing import Dict, Any, Optional

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# XML namespaces
NAMESPACES = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

# Register namespaces for lxml
for prefix, uri in NAMESPACES.items():
    etree.register_namespace(prefix, uri)


def hex_to_xml_color(hex_color: str) -> str:
    """Convert #RRGGBB to RRGGBB for XML."""
    return hex_color.lstrip('#').upper()


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert #RRGGBB to RGBColor."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


def add_design_elements_via_xml(extract_dir: Path, template_id: str, colors: Dict[str, str]) -> None:
    """
    Add visual design elements to slide master by manipulating XML directly.

    python-pptx doesn't support adding shapes to masters/layouts, so we
    inject rectangle shapes directly into the slideMaster1.xml file.
    """
    # Get colors with defaults
    primary = hex_to_xml_color(colors.get('primary', '#2c5aa0'))
    secondary = hex_to_xml_color(colors.get('secondary', '#4a4a4a'))
    accent = hex_to_xml_color(colors.get('accent', '#28a745'))
    surface = hex_to_xml_color(colors.get('surface', '#f8f9fa'))
    background = hex_to_xml_color(colors.get('background', '#ffffff'))

    slide_master_path = extract_dir / "ppt" / "slideMasters" / "slideMaster1.xml"
    if not slide_master_path.exists():
        logger.warning(f"Slide master not found: {slide_master_path}")
        return

    tree = etree.parse(str(slide_master_path))
    root = tree.getroot()

    # Find the shape tree (spTree)
    sp_tree = root.find('.//p:spTree', NAMESPACES)
    if sp_tree is None:
        logger.warning("Could not find spTree in slide master")
        return

    # Get design shapes for this template
    if template_id == 'corporate':
        shapes = _get_corporate_shapes_xml(primary, secondary, accent, surface)
    elif template_id == 'dark_mode':
        shapes = _get_dark_mode_shapes_xml(primary, secondary, accent, background)
    elif template_id == 'creative':
        shapes = _get_creative_shapes_xml(primary, secondary, accent)
    elif template_id == 'minimal':
        shapes = _get_minimal_shapes_xml(primary, secondary)
    elif template_id == 'educational':
        shapes = _get_educational_shapes_xml(primary, secondary, accent, surface)
    else:
        shapes = _get_corporate_shapes_xml(primary, secondary, accent, surface)

    # Insert shapes at the beginning of spTree (so they're behind other elements)
    for i, shape_xml in enumerate(shapes):
        shape_elem = etree.fromstring(shape_xml)
        sp_tree.insert(i + 2, shape_elem)  # After nvGrpSpPr and grpSpPr

    tree.write(str(slide_master_path), xml_declaration=True, encoding='UTF-8', standalone=True)
    logger.debug(f"Added design shapes to slide master for {template_id}")


def _make_rect_shape_xml(shape_id: int, name: str, x_emu: int, y_emu: int,
                          cx_emu: int, cy_emu: int, fill_color: str) -> str:
    """Generate XML for a rectangle shape with solid fill."""
    return f'''<p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="{name}"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{x_emu}" y="{y_emu}"/>
      <a:ext cx="{cx_emu}" cy="{cy_emu}"/>
    </a:xfrm>
    <a:prstGeom prst="rect">
      <a:avLst/>
    </a:prstGeom>
    <a:solidFill>
      <a:srgbClr val="{fill_color}"/>
    </a:solidFill>
    <a:ln>
      <a:noFill/>
    </a:ln>
  </p:spPr>
</p:sp>'''


# EMU conversion helpers (914400 EMU = 1 inch)
EMU_PER_INCH = 914400

def _inches_to_emu(inches: float) -> int:
    return int(inches * EMU_PER_INCH)


def _get_corporate_shapes_xml(primary: str, secondary: str, accent: str, surface: str) -> list:
    """Corporate: Blue header bar, green accent line, gray footer."""
    return [
        # Header bar - top
        _make_rect_shape_xml(
            100, "Header Bar",
            0, 0,
            _inches_to_emu(10), _inches_to_emu(0.4),
            primary
        ),
        # Accent line below title
        _make_rect_shape_xml(
            101, "Accent Line",
            _inches_to_emu(0.5), _inches_to_emu(1.5),
            _inches_to_emu(9), _inches_to_emu(0.03),
            accent
        ),
        # Footer bar
        _make_rect_shape_xml(
            102, "Footer Bar",
            0, _inches_to_emu(7.2),
            _inches_to_emu(10), _inches_to_emu(0.3),
            surface
        ),
    ]


def _get_dark_mode_shapes_xml(primary: str, secondary: str, accent: str, background: str) -> list:
    """Dark mode: Dark background with purple accent."""
    return [
        # Full dark background
        _make_rect_shape_xml(
            100, "Dark Background",
            0, 0,
            _inches_to_emu(10), _inches_to_emu(7.5),
            background
        ),
        # Top accent bar
        _make_rect_shape_xml(
            101, "Top Accent",
            0, 0,
            _inches_to_emu(10), _inches_to_emu(0.15),
            primary
        ),
        # Side accent
        _make_rect_shape_xml(
            102, "Side Accent",
            0, 0,
            _inches_to_emu(0.1), _inches_to_emu(7.5),
            secondary
        ),
    ]


def _get_creative_shapes_xml(primary: str, secondary: str, accent: str) -> list:
    """Creative: Bold orange header, corner accent."""
    return [
        # Large header bar
        _make_rect_shape_xml(
            100, "Header Bar",
            0, 0,
            _inches_to_emu(10), _inches_to_emu(0.5),
            primary
        ),
        # Bottom accent bar
        _make_rect_shape_xml(
            101, "Bottom Bar",
            0, _inches_to_emu(7.35),
            _inches_to_emu(10), _inches_to_emu(0.15),
            secondary
        ),
        # Corner accent (rectangle instead of triangle for simplicity)
        _make_rect_shape_xml(
            102, "Corner Accent",
            _inches_to_emu(9), _inches_to_emu(6.5),
            _inches_to_emu(1), _inches_to_emu(1),
            accent
        ),
    ]


def _get_minimal_shapes_xml(primary: str, secondary: str) -> list:
    """Minimal: Subtle lines only."""
    return [
        # Top line
        _make_rect_shape_xml(
            100, "Top Line",
            _inches_to_emu(0.5), _inches_to_emu(0.3),
            _inches_to_emu(9), _inches_to_emu(0.02),
            secondary
        ),
        # Bottom line
        _make_rect_shape_xml(
            101, "Bottom Line",
            _inches_to_emu(0.5), _inches_to_emu(7.2),
            _inches_to_emu(9), _inches_to_emu(0.02),
            secondary
        ),
    ]


def _get_educational_shapes_xml(primary: str, secondary: str, accent: str, surface: str) -> list:
    """Educational: Warm background with purple accents."""
    return [
        # Warm background tint
        _make_rect_shape_xml(
            100, "Background Tint",
            0, 0,
            _inches_to_emu(10), _inches_to_emu(7.5),
            surface
        ),
        # Header bar
        _make_rect_shape_xml(
            101, "Header Bar",
            0, 0,
            _inches_to_emu(10), _inches_to_emu(0.35),
            primary
        ),
        # Bottom accent
        _make_rect_shape_xml(
            102, "Bottom Accent",
            0, _inches_to_emu(7.25),
            _inches_to_emu(10), _inches_to_emu(0.25),
            accent
        ),
    ]


def create_base_presentation() -> Presentation:
    """Create a base presentation with standard slide layouts."""
    prs = Presentation()

    # Add a single blank slide to ensure all layouts are properly initialized
    # This helps python-pptx create all standard layouts
    blank_layout = prs.slide_layouts[6]  # Blank layout
    prs.slides.add_slide(blank_layout)

    return prs


def extract_pptx(pptx_path: Path, extract_dir: Path) -> None:
    """Extract PPTX (ZIP) contents to directory."""
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        zf.extractall(extract_dir)


def repackage_pptx(source_dir: Path, output_path: Path) -> None:
    """Repackage directory contents as PPTX (ZIP)."""
    # Remove existing file if present
    if output_path.exists():
        output_path.unlink()

    with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
        for root, dirs, files in os.walk(source_dir):
            for file in files:
                file_path = Path(root) / file
                arcname = file_path.relative_to(source_dir)
                zf.write(file_path, arcname)


def modify_theme_colors(theme_path: Path, colors: Dict[str, str], template_name: str) -> None:
    """
    Modify the theme XML to use custom colors.

    Color mapping from registry.json to PowerPoint theme slots:
    - primary -> accent1 (main brand color, used for shapes/charts)
    - secondary -> accent2
    - accent -> accent3
    - background -> lt1 (light 1 - slide background)
    - surface -> lt2 (light 2 - secondary background)
    - text_primary -> dk1 (dark 1 - main text)
    - text_secondary -> dk2 (dark 2 - secondary text)
    """
    # Parse the theme XML
    tree = etree.parse(str(theme_path))
    root = tree.getroot()

    # Find the color scheme element
    clr_scheme = root.find('.//a:clrScheme', NAMESPACES)
    if clr_scheme is None:
        logger.error(f"Could not find clrScheme in {theme_path}")
        return

    # Update the scheme name
    clr_scheme.set('name', template_name.replace('_', ' ').title())

    # Color mapping: registry key -> (theme element, is_srgb)
    # srgbClr for solid colors, sysClr for system colors (we'll replace with srgbClr)
    color_mapping = {
        'primary': 'accent1',
        'secondary': 'accent2',
        'accent': 'accent3',
        'text_primary': 'dk1',
        'text_secondary': 'dk2',
        'background': 'lt1',
        'surface': 'lt2',
    }

    for registry_key, theme_element in color_mapping.items():
        if registry_key not in colors:
            continue

        color_value = hex_to_xml_color(colors[registry_key])

        # Find the element (e.g., <a:accent1>)
        elem = clr_scheme.find(f'a:{theme_element}', NAMESPACES)
        if elem is None:
            logger.warning(f"Could not find {theme_element} in theme")
            continue

        # Clear existing color children
        for child in list(elem):
            elem.remove(child)

        # Add new srgbClr element
        srgb_elem = etree.SubElement(elem, f'{{{NAMESPACES["a"]}}}srgbClr')
        srgb_elem.set('val', color_value)

        logger.debug(f"Set {theme_element} to {color_value}")

    # Also set accent4-6 to variations of existing colors for consistency
    # accent4 = lighter version of primary
    # accent5 = lighter version of secondary
    # accent6 = lighter version of accent
    if 'primary' in colors:
        _set_accent_color(clr_scheme, 'accent4', colors.get('secondary', colors['primary']))
    if 'secondary' in colors:
        _set_accent_color(clr_scheme, 'accent5', colors.get('accent', colors['secondary']))
    if 'accent' in colors:
        _set_accent_color(clr_scheme, 'accent6', colors.get('primary', colors['accent']))

    # Write the modified XML back
    tree.write(str(theme_path), xml_declaration=True, encoding='UTF-8', standalone=True)
    logger.info(f"Updated theme colors for {template_name}")


def _set_accent_color(clr_scheme, element_name: str, hex_color: str) -> None:
    """Helper to set an accent color element."""
    elem = clr_scheme.find(f'a:{element_name}', NAMESPACES)
    if elem is None:
        return

    color_value = hex_to_xml_color(hex_color)

    for child in list(elem):
        elem.remove(child)

    srgb_elem = etree.SubElement(elem, f'{{{NAMESPACES["a"]}}}srgbClr')
    srgb_elem.set('val', color_value)


def create_themed_template(
    template_id: str,
    colors: Dict[str, str],
    output_path: Path,
    template_name: str
) -> bool:
    """
    Create a themed PPTX template.

    Args:
        template_id: Template identifier (e.g., "corporate")
        colors: Color dictionary from registry
        output_path: Where to save the template
        template_name: Display name for the template

    Returns:
        True if successful, False otherwise
    """
    logger.info(f"Creating template: {template_id}")

    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        base_pptx = temp_path / "base.pptx"
        extract_dir = temp_path / "extracted"

        # Step 1: Create base presentation
        prs = create_base_presentation()
        prs.save(str(base_pptx))
        logger.debug(f"Created base presentation at {base_pptx}")

        # Step 2: Extract PPTX contents
        extract_dir.mkdir()
        extract_pptx(base_pptx, extract_dir)
        logger.debug(f"Extracted to {extract_dir}")

        # Step 3: Modify theme colors
        theme_path = extract_dir / "ppt" / "theme" / "theme1.xml"
        if not theme_path.exists():
            logger.error(f"Theme file not found: {theme_path}")
            return False

        modify_theme_colors(theme_path, colors, template_name)

        # Step 4: Add visual design elements via XML manipulation
        add_design_elements_via_xml(extract_dir, template_id, colors)

        # Step 4: Repackage as PPTX
        output_path.parent.mkdir(parents=True, exist_ok=True)
        repackage_pptx(extract_dir, output_path)
        logger.info(f"Created template: {output_path}")

        return True


def load_registry(registry_path: Path) -> Dict[str, Any]:
    """Load the template registry."""
    with open(registry_path, 'r', encoding='utf-8') as f:
        return json.load(f)


def main():
    parser = argparse.ArgumentParser(
        description="Generate themed PPTX templates"
    )
    parser.add_argument(
        '--template', '-t',
        help="Generate only this template (default: all)"
    )
    parser.add_argument(
        '--verbose', '-v',
        action='store_true',
        help="Enable verbose output"
    )

    args = parser.parse_args()

    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)

    # Find paths
    script_dir = Path(__file__).parent
    templates_dir = script_dir.parent.parent / "templates" / "pptx"
    masters_dir = templates_dir / "masters"
    registry_path = masters_dir / "registry.json"

    if not registry_path.exists():
        logger.error(f"Registry not found: {registry_path}")
        return 1

    # Load registry
    registry = load_registry(registry_path)
    templates = registry.get('templates', {})

    if not templates:
        logger.error("No templates found in registry")
        return 1

    # Filter templates if specific one requested
    if args.template:
        if args.template not in templates:
            logger.error(f"Template '{args.template}' not found in registry")
            logger.info(f"Available: {', '.join(templates.keys())}")
            return 1
        templates = {args.template: templates[args.template]}

    # Generate templates
    success_count = 0
    for template_id, config in templates.items():
        output_path = masters_dir / config.get('path', f'{template_id}.pptx')
        colors = config.get('colors', {})
        name = config.get('name', template_id)

        if create_themed_template(template_id, colors, output_path, name):
            success_count += 1
        else:
            logger.error(f"Failed to create template: {template_id}")

    logger.info(f"Generated {success_count}/{len(templates)} templates")

    return 0 if success_count == len(templates) else 1


if __name__ == "__main__":
    exit(main())
