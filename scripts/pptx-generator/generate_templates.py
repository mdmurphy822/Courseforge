#!/usr/bin/env python3
"""
Generate Master Templates - Creates styled PPTX templates programmatically.

This script generates the built-in master templates for the Slideforge system.
Since python-pptx cannot modify master slides directly, we create presentations
with properly styled placeholders that serve as templates.

Usage:
    python generate_templates.py
"""

import json
import logging
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logging.basicConfig(level=logging.INFO, format='%(levelname)s: %(message)s')
logger = logging.getLogger(__name__)


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


def create_template(template_id: str, config: dict, output_dir: Path) -> Path:
    """
    Create a PPTX template with the specified configuration.

    Args:
        template_id: Template identifier (e.g., "corporate")
        config: Template configuration from registry
        output_dir: Directory to save the template

    Returns:
        Path to the created template file
    """
    prs = Presentation()
    colors = config.get("colors", {})
    fonts = config.get("fonts", {})

    primary = hex_to_rgb(colors.get("primary", "#2c5aa0"))
    secondary = hex_to_rgb(colors.get("secondary", "#4a4a4a"))
    accent = hex_to_rgb(colors.get("accent", "#28a745"))
    background = hex_to_rgb(colors.get("background", "#ffffff"))
    text_primary = hex_to_rgb(colors.get("text_primary", "#333333"))
    surface = hex_to_rgb(colors.get("surface", "#f8f9fa"))

    heading_font = fonts.get("heading", "Arial")
    body_font = fonts.get("body", "Arial")

    # Get slide layouts
    layouts = prs.slide_layouts

    # Create a sample slide of each type to "prime" the template
    # This ensures the layouts are properly initialized

    # 0: Title Slide
    slide = prs.slides.add_slide(layouts[0])
    if slide.shapes.title:
        slide.shapes.title.text = "Presentation Title"
        for para in slide.shapes.title.text_frame.paragraphs:
            para.font.name = heading_font
            para.font.size = Pt(44)
            para.font.color.rgb = primary
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = "Subtitle goes here"
        for para in subtitle.text_frame.paragraphs:
            para.font.name = body_font
            para.font.size = Pt(24)
            para.font.color.rgb = text_primary

    # 1: Title and Content
    slide = prs.slides.add_slide(layouts[1])
    if slide.shapes.title:
        slide.shapes.title.text = "Content Slide Title"
        for para in slide.shapes.title.text_frame.paragraphs:
            para.font.name = heading_font
            para.font.size = Pt(32)
            para.font.color.rgb = primary
    # Find content placeholder
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            tf = shape.text_frame
            tf.paragraphs[0].text = "First bullet point"
            p = tf.add_paragraph()
            p.text = "Second bullet point"
            p = tf.add_paragraph()
            p.text = "Third bullet point"
            for para in tf.paragraphs:
                para.font.name = body_font
                para.font.size = Pt(24)
                para.font.color.rgb = text_primary
            break

    # 2: Section Header
    slide = prs.slides.add_slide(layouts[2])
    if slide.shapes.title:
        slide.shapes.title.text = "Section Title"
        for para in slide.shapes.title.text_frame.paragraphs:
            para.font.name = heading_font
            para.font.size = Pt(40)
            para.font.color.rgb = primary

    # 3: Two Content
    slide = prs.slides.add_slide(layouts[3])
    if slide.shapes.title:
        slide.shapes.title.text = "Two Column Layout"
        for para in slide.shapes.title.text_frame.paragraphs:
            para.font.name = heading_font
            para.font.size = Pt(32)
            para.font.color.rgb = primary

    # 4: Comparison
    slide = prs.slides.add_slide(layouts[4])
    if slide.shapes.title:
        slide.shapes.title.text = "Comparison Layout"
        for para in slide.shapes.title.text_frame.paragraphs:
            para.font.name = heading_font
            para.font.size = Pt(32)
            para.font.color.rgb = primary

    # 5: Title Only
    slide = prs.slides.add_slide(layouts[5])
    if slide.shapes.title:
        slide.shapes.title.text = "Title Only Layout"
        for para in slide.shapes.title.text_frame.paragraphs:
            para.font.name = heading_font
            para.font.size = Pt(32)
            para.font.color.rgb = primary

    # 6: Blank
    slide = prs.slides.add_slide(layouts[6])

    # Save template
    output_path = output_dir / f"{template_id}.pptx"
    prs.save(str(output_path))
    logger.info(f"Created template: {output_path}")

    return output_path


def main():
    """Generate all master templates from registry."""
    script_dir = Path(__file__).parent
    templates_dir = script_dir.parent.parent / "templates" / "pptx"
    masters_dir = templates_dir / "masters"
    registry_path = masters_dir / "registry.json"

    if not registry_path.exists():
        logger.error(f"Registry not found: {registry_path}")
        return

    with open(registry_path, 'r') as f:
        registry = json.load(f)

    templates = registry.get("templates", {})

    logger.info(f"Generating {len(templates)} templates...")

    for template_id, config in templates.items():
        try:
            create_template(template_id, config, masters_dir)
        except Exception as e:
            logger.error(f"Failed to create {template_id}: {e}")

    logger.info("Template generation complete!")

    # List created templates
    print("\n=== Generated Templates ===")
    for pptx_file in masters_dir.glob("*.pptx"):
        print(f"  {pptx_file.name}")


if __name__ == "__main__":
    main()
