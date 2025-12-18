"""Inspect commands for Slideforge CLI."""

import json
from collections import Counter
from pathlib import Path

import click


@click.group('inspect')
def inspect_cmd():
    """Inspect presentation files and show statistics."""
    pass


@inspect_cmd.command('presentation')
@click.argument('file', type=click.Path(exists=True))
@click.option('--json-output', '-j', is_flag=True, help='Output as JSON')
def inspect_presentation(file, json_output):
    """Show presentation statistics.

    FILE: Path to presentation JSON file
    """
    file_path = Path(file)

    try:
        with open(file_path) as f:
            data = json.load(f)
    except json.JSONDecodeError as e:
        click.echo(f"Error: Invalid JSON - {e}", err=True)
        raise SystemExit(1)

    # Gather statistics
    metadata = data.get('metadata', {})
    sections = data.get('sections', [])

    stats = {
        'title': metadata.get('title', 'Untitled'),
        'author': metadata.get('author', 'Unknown'),
        'section_count': len(sections),
        'slide_count': 0,
        'slide_types': Counter(),
        'slides_with_notes': 0,
        'total_bullets': 0,
        'six_six_violations': 0,
        'sections': []
    }

    for section in sections:
        section_title = section.get('title', 'Untitled Section')
        slides = section.get('slides', [])
        stats['slide_count'] += len(slides)

        section_stats = {
            'title': section_title,
            'slide_count': len(slides),
            'types': Counter()
        }

        for slide in slides:
            slide_type = slide.get('type', 'content')
            stats['slide_types'][slide_type] += 1
            section_stats['types'][slide_type] += 1

            if slide.get('notes'):
                stats['slides_with_notes'] += 1

            content = slide.get('content', {})

            # Count bullets
            for key in ['bullets', 'left', 'right']:
                items = content.get(key, [])
                stats['total_bullets'] += len(items)
                if len(items) > 6:
                    stats['six_six_violations'] += 1
                for item in items:
                    if len(item.split()) > 6:
                        stats['six_six_violations'] += 1

        section_stats['types'] = dict(section_stats['types'])
        stats['sections'].append(section_stats)

    stats['slide_types'] = dict(stats['slide_types'])
    stats['notes_coverage'] = round(
        (stats['slides_with_notes'] / stats['slide_count'] * 100)
        if stats['slide_count'] > 0 else 0, 1
    )

    if json_output:
        click.echo(json.dumps(stats, indent=2))
        return

    # Pretty print
    click.echo(f"\nPresentation: {stats['title']}")
    click.echo("=" * 50)
    click.echo(f"Author: {stats['author']}")
    click.echo(f"Sections: {stats['section_count']}")
    click.echo(f"Total slides: {stats['slide_count']}")
    click.echo(f"Notes coverage: {stats['notes_coverage']}%")
    click.echo(f"6x6 violations: {stats['six_six_violations']}")

    click.echo("\nSlide Types:")
    for stype, count in sorted(stats['slide_types'].items()):
        click.echo(f"  {stype}: {count}")

    click.echo("\nSections:")
    for section in stats['sections']:
        click.echo(f"  {section['title']}: {section['slide_count']} slides")


@inspect_cmd.command('pptx')
@click.argument('file', type=click.Path(exists=True))
@click.option('--json-output', '-j', is_flag=True, help='Output as JSON')
def inspect_pptx(file, json_output):
    """Analyze existing PPTX file.

    FILE: Path to PPTX file
    """
    file_path = Path(file)

    try:
        from pptx import Presentation
    except ImportError:
        click.echo("Error: python-pptx not installed", err=True)
        raise SystemExit(1)

    try:
        prs = Presentation(str(file_path))
    except Exception as e:
        click.echo(f"Error: Could not open PPTX - {e}", err=True)
        raise SystemExit(1)

    stats = {
        'filename': file_path.name,
        'slide_count': len(prs.slides),
        'slide_width': prs.slide_width.inches,
        'slide_height': prs.slide_height.inches,
        'aspect_ratio': f"{prs.slide_width.inches / prs.slide_height.inches:.2f}:1",
        'slides_with_notes': 0,
        'layouts_used': Counter()
    }

    for slide in prs.slides:
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
        stats['layouts_used'][layout_name] += 1

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            stats['slides_with_notes'] += 1

    stats['layouts_used'] = dict(stats['layouts_used'])
    stats['notes_coverage'] = round(
        (stats['slides_with_notes'] / stats['slide_count'] * 100)
        if stats['slide_count'] > 0 else 0, 1
    )

    if json_output:
        click.echo(json.dumps(stats, indent=2))
        return

    click.echo(f"\nPPTX Analysis: {stats['filename']}")
    click.echo("=" * 50)
    click.echo(f"Slides: {stats['slide_count']}")
    click.echo(f"Dimensions: {stats['slide_width']}\" x {stats['slide_height']}\"")
    click.echo(f"Aspect ratio: {stats['aspect_ratio']}")
    click.echo(f"Notes coverage: {stats['notes_coverage']}%")

    click.echo("\nLayouts used:")
    for layout, count in sorted(stats['layouts_used'].items()):
        click.echo(f"  {layout}: {count}")


@inspect_cmd.command('theme')
@click.argument('file', type=click.Path(exists=True))
def inspect_theme(file):
    """Show theme configuration details.

    FILE: Path to theme JSON file
    """
    file_path = Path(file)

    try:
        with open(file_path) as f:
            data = json.load(f)
    except json.JSONDecodeError as e:
        click.echo(f"Error: Invalid JSON - {e}", err=True)
        raise SystemExit(1)

    click.echo(f"\nTheme: {data.get('name', file_path.stem)}")
    click.echo("=" * 50)

    if 'description' in data:
        click.echo(f"Description: {data['description']}")

    if 'aspect_ratio' in data:
        click.echo(f"Aspect ratio: {data['aspect_ratio']}")

    colors = data.get('colors', {})
    if colors:
        click.echo("\nColors:")
        for name, value in colors.items():
            click.echo(f"  {name}: {value}")

    fonts = data.get('fonts', {})
    if fonts:
        click.echo("\nFonts:")
        if 'heading' in fonts:
            click.echo(f"  Heading: {fonts['heading']}")
        if 'body' in fonts:
            click.echo(f"  Body: {fonts['body']}")

    sizes = data.get('sizes', {})
    if sizes:
        click.echo("\nFont sizes:")
        for name, size in sizes.items():
            click.echo(f"  {name}: {size}pt")
