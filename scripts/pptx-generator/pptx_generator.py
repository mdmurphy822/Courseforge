#!/usr/bin/env python3
"""
PPTX Generator - Create PowerPoint presentations from structured content.

This module provides the core functionality for generating PPTX files from
structured JSON content using the python-pptx library.

Features:
- Template-based theming (optional .pptx template)
- Multiple slide layout types
- Automatic content fitting
- Speaker notes support
- Image and shape integration
- Consistent styling

Usage:
    from pptx_generator import PPTXGenerator

    generator = PPTXGenerator(template_path="theme.pptx")
    generator.create_presentation(structured_content)
    generator.save("output.pptx")
"""

import json
import logging
import os
import sys
from dataclasses import dataclass, field
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# Import template loader (optional - graceful fallback if not available)
try:
    from template_loader import TemplateLoader, TemplateConfig, TemplateColors as LoaderColors
    TEMPLATE_LOADER_AVAILABLE = True
except ImportError:
    TEMPLATE_LOADER_AVAILABLE = False

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


class SlideType(Enum):
    """Available slide types"""
    TITLE = "title"
    SECTION_HEADER = "section_header"
    CONTENT = "content"
    TWO_CONTENT = "two_content"
    COMPARISON = "comparison"
    TITLE_ONLY = "title_only"
    BLANK = "blank"
    IMAGE = "image"
    QUOTE = "quote"
    TABLE = "table"
    PROCESS_FLOW = "process_flow"
    COMPARISON_MATRIX = "comparison_matrix"
    TIMELINE = "timeline"
    KEY_VALUE = "key_value"
    CALLOUT = "callout"
    STATS_GRID = "stats_grid"
    CARDS_GRID = "cards_grid"
    KEY_CONCEPT = "key_concept"
    DATA_VIZ = "data_viz"
    DIVIDER = "divider"
    AGENDA = "agenda"
    IMAGE_GRID = "image_grid"


@dataclass
class SlideContent:
    """Represents content for a single slide"""
    slide_type: SlideType
    title: str = ""
    subtitle: str = ""
    bullets: List[str] = field(default_factory=list)
    left_content: List[str] = field(default_factory=list)
    right_content: List[str] = field(default_factory=list)
    image_path: Optional[str] = None
    image_alt: str = ""
    quote_text: str = ""
    quote_attribution: str = ""
    notes: str = ""


@dataclass
class PresentationMetadata:
    """Presentation metadata"""
    title: str = "Untitled Presentation"
    subtitle: str = ""
    author: str = ""
    date: str = ""
    subject: str = ""


@dataclass
class ThemeColors:
    """Theme color configuration with full semantic color support"""
    primary: str = "#2c5aa0"
    secondary: str = "#28a745"
    accent: str = "#ffc107"
    background: str = "#ffffff"
    surface: str = "#f8f9fa"          # For cards, containers, alt rows
    text_dark: str = "#333333"
    text_light: str = "#ffffff"
    text_secondary: str = "#666666"   # For secondary/muted text
    success: str = "#28a745"          # Semantic: positive indicators
    warning: str = "#ffc107"          # Semantic: caution indicators
    error: str = "#dc3545"            # Semantic: negative indicators
    info: str = "#17a2b8"             # Semantic: informational


class PPTXGenerator:
    """
    PowerPoint presentation generator.

    Creates PPTX files from structured content with support for
    templates, themes, and multiple slide layouts.
    """

    # Default slide layout indices (standard PowerPoint)
    LAYOUT_INDICES = {
        SlideType.TITLE: 0,           # Title Slide
        SlideType.CONTENT: 1,         # Title and Content
        SlideType.SECTION_HEADER: 2,  # Section Header
        SlideType.TWO_CONTENT: 3,     # Two Content
        SlideType.COMPARISON: 4,      # Comparison
        SlideType.TITLE_ONLY: 5,      # Title Only
        SlideType.BLANK: 6,           # Blank
        SlideType.IMAGE: 1,           # Use content layout for images
        SlideType.QUOTE: 5,           # Use title only for quotes
        SlideType.TABLE: 5,           # Use title only for tables
        SlideType.PROCESS_FLOW: 5,    # Use title only for process flows
        SlideType.COMPARISON_MATRIX: 5,  # Use title only for comparison matrix
        SlideType.TIMELINE: 5,        # Use title only for timelines
        SlideType.KEY_VALUE: 5,       # Use title only for key-value pairs
        SlideType.CALLOUT: 5,         # Custom layout for callout slides
        SlideType.STATS_GRID: 5,      # Custom layout for stats grid
        SlideType.CARDS_GRID: 5,      # Custom layout for cards grid
        SlideType.KEY_CONCEPT: 5,     # Custom layout for key concept
        SlideType.AGENDA: 5,          # Use title only for agenda slides
        SlideType.IMAGE_GRID: 5,      # Use title only for image grid
    }

    def __init__(
        self,
        template_path: Optional[str] = None,
        template_name: Optional[str] = None,
        theme_colors: Optional[ThemeColors] = None
    ):
        """
        Initialize the PPTX generator.

        Args:
            template_path: Path to .pptx template file for theming
            template_name: Name of registered template (e.g., "corporate", "dark_mode")
            theme_colors: Custom color scheme (overrides template colors)
        """
        self._template_config: Optional[TemplateConfig] = None
        self._template_loader: Optional[TemplateLoader] = None
        self._layout_overrides: Dict[str, int] = {}

        # Try to use TemplateLoader for registered templates
        if TEMPLATE_LOADER_AVAILABLE and template_name:
            self._template_loader = TemplateLoader()
            resolved_path, config = self._template_loader.get_template(template_name)
            self._template_config = config

            if resolved_path and resolved_path.exists():
                self.prs = Presentation(str(resolved_path))
                logger.info(f"Loaded template '{template_name}': {resolved_path}")
                # Use template's layout mappings
                self._layout_overrides = config.layouts
                # Use template colors if no override provided
                if theme_colors is None:
                    theme_colors = ThemeColors(
                        primary=config.colors.primary,
                        secondary=config.colors.secondary,
                        accent=config.colors.accent,
                        background=config.colors.background,
                        surface=config.colors.surface,
                        text_dark=config.colors.text_primary,
                        text_light="#ffffff",
                        text_secondary=config.colors.text_secondary,
                        success=config.colors.success,
                        warning=config.colors.warning,
                        error=getattr(config.colors, 'error', '#dc3545'),
                        info=config.colors.info,
                    )
            else:
                logger.warning(f"Template '{template_name}' not found, using default")
                self.prs = Presentation()

        # Fall back to direct path
        elif template_path and Path(template_path).exists():
            self.prs = Presentation(template_path)
            logger.info(f"Loaded template: {template_path}")

        else:
            self.prs = Presentation()
            logger.info("Created new presentation (no template)")

        self.theme_colors = theme_colors or ThemeColors()
        self.metadata = PresentationMetadata()
        self._slide_count = 0
        # Cache container colors derived from theme
        self._container_colors = self._build_container_colors()

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color string to RGBColor object.

        Args:
            hex_color: Color in hex format (e.g., "#2c5aa0" or "2c5aa0")

        Returns:
            RGBColor object for use with python-pptx
        """
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )

    def _lighten_color(self, hex_color: str, factor: float = 0.85) -> str:
        """Create a lighter version of a color for backgrounds.

        Args:
            hex_color: Base color in hex format
            factor: Lightening factor (0.0 = original, 1.0 = white)

        Returns:
            Lightened color as hex string
        """
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)

        # Blend toward white
        r = int(r + (255 - r) * factor)
        g = int(g + (255 - g) * factor)
        b = int(b + (255 - b) * factor)

        return f"#{r:02x}{g:02x}{b:02x}"

    def _build_container_colors(self) -> Dict[str, Dict[str, str]]:
        """Build container/callout colors derived from theme semantic colors.

        Returns:
            Dictionary mapping container types to their bg/text colors
        """
        return {
            "info": {
                "bg": self._lighten_color(self.theme_colors.info, 0.85),
                "text": self.theme_colors.info
            },
            "success": {
                "bg": self._lighten_color(self.theme_colors.success, 0.85),
                "text": self.theme_colors.success
            },
            "warning": {
                "bg": self._lighten_color(self.theme_colors.warning, 0.85),
                "text": self._darken_color(self.theme_colors.warning, 0.4)
            },
            "tip": {
                "bg": self._lighten_color(self.theme_colors.accent, 0.85),
                "text": self._darken_color(self.theme_colors.accent, 0.4)
            },
            "error": {
                "bg": self._lighten_color(self.theme_colors.error, 0.85),
                "text": self.theme_colors.error
            },
        }

    def _darken_color(self, hex_color: str, factor: float = 0.3) -> str:
        """Create a darker version of a color for text on light backgrounds.

        Args:
            hex_color: Base color in hex format
            factor: Darkening factor (0.0 = original, 1.0 = black)

        Returns:
            Darkened color as hex string
        """
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)

        # Blend toward black
        r = int(r * (1 - factor))
        g = int(g * (1 - factor))
        b = int(b * (1 - factor))

        return f"#{r:02x}{g:02x}{b:02x}"

    def set_metadata(self, metadata: PresentationMetadata) -> None:
        """Set presentation metadata."""
        self.metadata = metadata

        # Set core properties
        core_props = self.prs.core_properties
        core_props.title = metadata.title
        core_props.author = metadata.author
        core_props.subject = metadata.subject

        logger.info(f"Set metadata - Title: {metadata.title}")

    def add_title_slide(
        self,
        title: str,
        subtitle: str = "",
        notes: str = ""
    ) -> None:
        """
        Add a title slide.

        Args:
            title: Main title text
            subtitle: Subtitle or author/date line
            notes: Speaker notes
        """
        layout = self._get_layout(SlideType.TITLE)
        slide = self.prs.slides.add_slide(layout)

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = title

        # Set subtitle
        if len(slide.placeholders) > 1:
            subtitle_placeholder = slide.placeholders[1]
            subtitle_placeholder.text = subtitle

        # Add speaker notes
        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added title slide: {title}")

    def add_section_header(
        self,
        title: str,
        subtitle: str = "",
        notes: str = ""
    ) -> None:
        """
        Add a section header slide.

        Args:
            title: Section title
            subtitle: Section description
            notes: Speaker notes
        """
        layout = self._get_layout(SlideType.SECTION_HEADER)
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added section header: {title}")


    def add_divider_slide(self, slide_data: Dict[str, Any]) -> None:
        """
        Add a divider slide for visual section transitions.

        Creates a visually striking divider slide to separate major presentation sections.
        Supports multiple visual styles for different presentation aesthetics.

        Args:
            slide_data: Dictionary containing divider configuration
                - title: Main divider title text
                - content: Dict with:
                    - section_number: Optional section number (e.g., "01", "02", "A")
                    - subtitle: Optional subtitle text
                    - divider_style: Style variant ("bold", "minimal", "graphic", "numbered")
                    - accent_color: Color theme ("primary", "secondary", "accent")

        Divider Styles:
            - bold: Large centered title with horizontal accent lines above/below
            - minimal: Left-aligned title with vertical accent bar
            - graphic: Large circle with title inside
            - numbered: Large faded section number background with title overlay
        """
        # Extract configuration
        title = slide_data.get("title", "")
        content = slide_data.get("content", {})
        notes = slide_data.get("notes", "")

        section_number = content.get("section_number", "")
        subtitle = content.get("subtitle", "")
        divider_style = content.get("divider_style", "bold")
        accent_color_key = content.get("accent_color", "primary")

        # Get accent color from theme
        color_map = {
            "primary": self.theme_colors.primary,
            "secondary": self.theme_colors.secondary,
            "accent": self.theme_colors.accent,
        }
        accent_color_hex = color_map.get(accent_color_key, self.theme_colors.primary)
        accent_color = RGBColor.from_string(accent_color_hex[1:])

        # Use blank layout for full control
        layout = self._get_layout(SlideType.BLANK)
        slide = self.prs.slides.add_slide(layout)

        # Render based on style
        if divider_style == "bold":
            self._render_bold_divider(slide, title, subtitle, section_number, accent_color)
        elif divider_style == "minimal":
            self._render_minimal_divider(slide, title, subtitle, section_number, accent_color)
        elif divider_style == "graphic":
            self._render_graphic_divider(slide, title, subtitle, section_number, accent_color)
        elif divider_style == "numbered":
            self._render_numbered_divider(slide, title, subtitle, section_number, accent_color)
        else:
            # Default to bold if unknown style
            self._render_bold_divider(slide, title, subtitle, section_number, accent_color)

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added divider slide: {title} (style: {divider_style})")

    def _render_bold_divider(
        self,
        slide,
        title: str,
        subtitle: str,
        section_number: str,
        accent_color: RGBColor
    ) -> None:
        """Render bold divider style: horizontal lines with centered title."""
        # Top horizontal line
        top_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(2.0),
            Inches(2.5),
            Inches(6.0),
            Inches(0.05)
        )
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = accent_color
        top_line.line.fill.background()

        # Section number (if provided)
        if section_number:
            number_box = slide.shapes.add_textbox(
                Inches(2.0),
                Inches(1.8),
                Inches(6.0),
                Inches(0.5)
            )
            tf = number_box.text_frame
            tf.paragraphs[0].text = section_number
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = accent_color
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Main title (large, centered between lines)
        title_box = slide.shapes.add_textbox(
            Inches(1.5),
            Inches(2.7),
            Inches(7.0),
            Inches(1.2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(48)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme_colors.text_dark)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Bottom horizontal line
        bottom_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(2.0),
            Inches(4.2),
            Inches(6.0),
            Inches(0.05)
        )
        bottom_line.fill.solid()
        bottom_line.fill.fore_color.rgb = accent_color
        bottom_line.line.fill.background()

        # Subtitle (if provided)
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(2.0),
                Inches(4.5),
                Inches(6.0),
                Inches(0.6)
            )
            tf = subtitle_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = subtitle
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme_colors.text_secondary)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _render_minimal_divider(
        self,
        slide,
        title: str,
        subtitle: str,
        section_number: str,
        accent_color: RGBColor
    ) -> None:
        """Render minimal divider style: vertical accent bar with left-aligned title."""
        # Vertical accent bar on left
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.0),
            Inches(2.0),
            Inches(0.1),
            Inches(3.5)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_color
        accent_bar.line.fill.background()

        # Section number (if provided)
        if section_number:
            number_box = slide.shapes.add_textbox(
                Inches(1.5),
                Inches(2.0),
                Inches(7.0),
                Inches(0.5)
            )
            tf = number_box.text_frame
            tf.paragraphs[0].text = section_number
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = accent_color
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        # Main title (left-aligned, large)
        title_top = 2.5 if section_number else 2.0
        title_box = slide.shapes.add_textbox(
            Inches(1.5),
            Inches(title_top),
            Inches(7.0),
            Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme_colors.text_dark)
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        # Subtitle (if provided)
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1.5),
                Inches(4.2),
                Inches(7.0),
                Inches(0.8)
            )
            tf = subtitle_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = subtitle
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme_colors.text_secondary)
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    def _render_graphic_divider(
        self,
        slide,
        title: str,
        subtitle: str,
        section_number: str,
        accent_color: RGBColor
    ) -> None:
        """Render graphic divider style: large circle with title inside."""
        # Large centered circle
        circle_size = 4.0
        circle_left = (10.0 - circle_size) / 2
        circle_top = (7.5 - circle_size) / 2

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(circle_left),
            Inches(circle_top),
            Inches(circle_size),
            Inches(circle_size)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent_color
        circle.line.fill.background()

        # Section number (if provided) - small text above title in circle
        title_offset = 0.0
        if section_number:
            title_offset = 0.3
            number_box = slide.shapes.add_textbox(
                Inches(circle_left + 0.2),
                Inches(circle_top + 1.2),
                Inches(circle_size - 0.4),
                Inches(0.4)
            )
            tf = number_box.text_frame
            tf.paragraphs[0].text = section_number
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme_colors.text_light)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Title inside circle (light text, centered)
        title_box = slide.shapes.add_textbox(
            Inches(circle_left + 0.3),
            Inches(circle_top + 1.5 + title_offset),
            Inches(circle_size - 0.6),
            Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme_colors.text_light)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Subtitle below circle (if provided)
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(2.0),
                Inches(circle_top + circle_size + 0.3),
                Inches(6.0),
                Inches(0.6)
            )
            tf = subtitle_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = subtitle
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme_colors.text_secondary)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _render_numbered_divider(
        self,
        slide,
        title: str,
        subtitle: str,
        section_number: str,
        accent_color: RGBColor
    ) -> None:
        """Render numbered divider style: large faded number with title overlay."""
        # Large faded section number background (right side)
        if section_number:
            number_bg = slide.shapes.add_textbox(
                Inches(5.0),
                Inches(1.5),
                Inches(4.5),
                Inches(4.5)
            )
            tf = number_bg.text_frame
            tf.paragraphs[0].text = section_number
            tf.paragraphs[0].font.size = Pt(96)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme_colors.surface)
            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
            tf.anchor = MSO_ANCHOR.MIDDLE

        # Main title (overlayed on left, themed text)
        title_box = slide.shapes.add_textbox(
            Inches(1.0),
            Inches(2.5),
            Inches(6.0),
            Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(44)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme_colors.text_dark)
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        # Accent line under title
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.0),
            Inches(4.2),
            Inches(3.0),
            Inches(0.08)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_color
        accent_line.line.fill.background()

        # Subtitle (if provided)
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1.0),
                Inches(4.5),
                Inches(6.0),
                Inches(0.8)
            )
            tf = subtitle_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = subtitle
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme_colors.text_secondary)
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    def add_content_slide(
        self,
        title: str,
        bullets: List[str],
        notes: str = ""
    ) -> None:
        """
        Add a content slide with bullet points.

        Args:
            title: Slide title
            bullets: List of bullet point strings
            notes: Speaker notes
        """
        layout = self._get_layout(SlideType.CONTENT)
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Find content placeholder
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Content placeholder
                content_placeholder = shape
                break

        if content_placeholder and bullets:
            tf = content_placeholder.text_frame
            tf.clear()

            for i, bullet in enumerate(bullets):
                if i == 0:
                    tf.paragraphs[0].text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added content slide: {title} ({len(bullets)} bullets)")

    def add_two_content_slide(
        self,
        title: str,
        left_content: List[str],
        right_content: List[str],
        notes: str = ""
    ) -> None:
        """
        Add a two-column content slide.

        Args:
            title: Slide title
            left_content: Left column bullet points
            right_content: Right column bullet points
            notes: Speaker notes
        """
        layout = self._get_layout(SlideType.TWO_CONTENT)
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Populate left and right placeholders
        placeholders = list(slide.placeholders)
        if len(placeholders) >= 3:
            self._fill_text_frame(placeholders[1].text_frame, left_content)
            self._fill_text_frame(placeholders[2].text_frame, right_content)

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added two-content slide: {title}")

    def add_comparison_slide(
        self,
        title: str,
        left_title: str,
        left_content: List[str],
        right_title: str,
        right_content: List[str],
        notes: str = ""
    ) -> None:
        """
        Add a comparison slide with two labeled columns.

        Args:
            title: Slide title
            left_title: Left column header
            left_content: Left column bullet points
            right_title: Right column header
            right_content: Right column bullet points
            notes: Speaker notes
        """
        layout = self._get_layout(SlideType.COMPARISON)
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # For comparison layout, we have title placeholders for each column
        placeholders = list(slide.placeholders)

        # Fill in column titles and content
        # Layout varies by template, so we handle this gracefully
        self._fill_comparison_content(slide, left_title, left_content,
                                       right_title, right_content)

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added comparison slide: {title}")

    def add_image_slide(
        self,
        title: str,
        image_path: str,
        alt_text: str = "",
        notes: str = ""
    ) -> None:
        """
        Add a slide with an image.

        Args:
            title: Slide title
            image_path: Path to image file
            alt_text: Alternative text for accessibility
            notes: Speaker notes
        """
        if not Path(image_path).exists():
            logger.warning(f"Image not found: {image_path}")
            return

        layout = self._get_layout(SlideType.TITLE_ONLY)
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Add image centered below title
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)

        pic = slide.shapes.add_picture(image_path, left, top, width=width)

        # Set alt text for accessibility
        pic.name = alt_text or title

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added image slide: {title}")

    def add_quote_slide(
        self,
        quote_text: str,
        attribution: str = "",
        notes: str = ""
    ) -> None:
        """
        Add a quote slide.

        Args:
            quote_text: The quote text
            attribution: Quote attribution (author, source)
            notes: Speaker notes
        """
        layout = self._get_layout(SlideType.BLANK)
        slide = self.prs.slides.add_slide(layout)

        # Add quote as centered text box
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(2)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        # Quote text
        p = tf.paragraphs[0]
        p.text = f'"{quote_text}"'
        p.font.size = Pt(28)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

        # Attribution
        if attribution:
            attr_p = tf.add_paragraph()
            attr_p.text = f"— {attribution}"
            attr_p.font.size = Pt(18)
            attr_p.alignment = PP_ALIGN.CENTER

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added quote slide")

    def add_blank_slide(self, notes: str = "") -> None:
        """Add a blank slide."""
        layout = self._get_layout(SlideType.BLANK)
        slide = self.prs.slides.add_slide(layout)

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info("Added blank slide")

    def add_table_slide(
        self,
        title: str,
        headers: List[str],
        rows: List[List[str]],
        totals_row: Optional[List[str]] = None,
        notes: str = ""
    ) -> None:
        """
        Add a slide with a data table.

        Args:
            title: Slide title
            headers: Column headers
            rows: Table data rows
            totals_row: Optional totals/summary row
            notes: Speaker notes
        """
        layout = self._get_layout(SlideType.TITLE_ONLY)
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Calculate table dimensions
        num_cols = len(headers)
        num_rows = len(rows) + 1  # +1 for header
        if totals_row:
            num_rows += 1

        # Table position and size
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(9.0)
        height = Inches(0.4 * num_rows)

        # Create table
        table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

        # Style settings - use theme colors
        header_color = self._hex_to_rgb(self.theme_colors.primary)
        alt_row_color = self._hex_to_rgb(self.theme_colors.surface)
        totals_row_color = self._hex_to_rgb(self._darken_color(self.theme_colors.surface, 0.1))

        # Set column widths (equal distribution)
        col_width = int(Inches(9.0 / num_cols))
        for col_idx in range(num_cols):
            table.columns[col_idx].width = col_width

        # Add headers
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = self._hex_to_rgb(self.theme_colors.text_light)
            paragraph.alignment = PP_ALIGN.CENTER

        # Add data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_text)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(16)
                paragraph.alignment = PP_ALIGN.LEFT

                # Alternate row coloring
                if row_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = alt_row_color

        # Add totals row if provided
        if totals_row:
            totals_idx = num_rows - 1
            for col_idx, cell_text in enumerate(totals_row):
                cell = table.cell(totals_idx, col_idx)
                cell.text = str(cell_text)
                cell.fill.solid()
                cell.fill.fore_color.rgb = totals_row_color
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.bold = True
                paragraph.font.size = Pt(16)

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added table slide: {title} ({num_cols}x{len(rows)})")

    def add_process_flow_slide(
        self,
        title: str,
        steps: List[Dict[str, str]],
        flow_style: str = "horizontal",
        notes: str = ""
    ) -> None:
        """
        Add a process flow slide showing sequential steps.

        Args:
            title: Slide title
            steps: List of step dicts with 'label' and optional 'description'
            flow_style: 'horizontal' or 'vertical'
            notes: Speaker notes
        """
        layout = self._get_layout(SlideType.TITLE_ONLY)
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        num_steps = min(len(steps), 6)  # Max 6 steps

        if flow_style == "horizontal":
            self._add_horizontal_process(slide, steps[:num_steps])
        else:
            self._add_vertical_process(slide, steps[:num_steps])

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added process flow slide: {title} ({num_steps} steps)")

    def _add_horizontal_process(self, slide, steps: List[Dict[str, str]]) -> None:
        """Add horizontal process flow."""
        num_steps = len(steps)
        total_width = 9.0
        step_width = total_width / num_steps
        arrow_width = 0.3
        circle_size = 0.8
        start_left = 0.5

        primary_color = RGBColor.from_string(self.theme_colors.primary[1:])

        for i, step in enumerate(steps):
            left = start_left + (i * step_width)

            # Step circle with number
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left + (step_width - circle_size) / 2),
                Inches(2.0),
                Inches(circle_size),
                Inches(circle_size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = primary_color
            circle.line.fill.background()

            # Number in circle
            tf = circle.text_frame
            tf.paragraphs[0].text = str(i + 1)
            tf.paragraphs[0].font.size = Pt(20)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE

            # Step label
            label_box = slide.shapes.add_textbox(
                Inches(left),
                Inches(3.0),
                Inches(step_width - 0.2),
                Inches(0.5)
            )
            tf = label_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = step.get("label", f"Step {i+1}")
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Step description
            if step.get("description"):
                desc_box = slide.shapes.add_textbox(
                    Inches(left),
                    Inches(3.5),
                    Inches(step_width - 0.2),
                    Inches(1.0)
                )
                tf = desc_box.text_frame
                tf.word_wrap = True
                tf.paragraphs[0].text = step["description"]
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Arrow to next step (except last)
            if i < num_steps - 1:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    Inches(left + step_width - arrow_width),
                    Inches(2.25),
                    Inches(arrow_width),
                    Inches(0.3)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = RGBColor(200, 200, 200)
                arrow.line.fill.background()

    def _add_vertical_process(self, slide, steps: List[Dict[str, str]]) -> None:
        """Add vertical process flow."""
        num_steps = len(steps)
        step_height = 0.9
        start_top = 1.8
        circle_size = 0.5

        primary_color = RGBColor.from_string(self.theme_colors.primary[1:])

        for i, step in enumerate(steps):
            top = start_top + (i * step_height)

            # Step circle with number
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.5),
                Inches(top),
                Inches(circle_size),
                Inches(circle_size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = primary_color
            circle.line.fill.background()

            # Number in circle
            tf = circle.text_frame
            tf.paragraphs[0].text = str(i + 1)
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE

            # Step text
            text_box = slide.shapes.add_textbox(
                Inches(1.2),
                Inches(top),
                Inches(8.0),
                Inches(step_height - 0.1)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = step.get("label", f"Step {i+1}")
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.bold = True

            if step.get("description"):
                p = tf.add_paragraph()
                p.text = step["description"]
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(102, 102, 102)

            # Connector line (except last)
            if i < num_steps - 1:
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.7),
                    Inches(top + circle_size),
                    Inches(0.1),
                    Inches(step_height - circle_size)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(200, 200, 200)
                line.line.fill.background()

    def add_comparison_matrix_slide(
        self,
        title: str,
        options: List[str],
        criteria: List[Dict[str, Any]],
        highlight_option: Optional[str] = None,
        notes: str = ""
    ) -> None:
        """
        Add a comparison matrix slide.

        Args:
            title: Slide title
            options: List of options to compare (columns)
            criteria: List of criteria dicts with 'name' and 'values'
            highlight_option: Option to highlight (column name)
            notes: Speaker notes
        """
        layout = self._get_layout(SlideType.TITLE_ONLY)
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        num_cols = len(options) + 1  # +1 for criteria column
        num_rows = len(criteria) + 1  # +1 for header

        # Table position
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(9.0)
        height = Inches(0.5 * num_rows)

        table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

        header_color = RGBColor.from_string(self.theme_colors.primary[1:])
        highlight_color = RGBColor.from_string(self.theme_colors.accent[1:])
        highlight_idx = options.index(highlight_option) + 1 if highlight_option in options else -1

        # Set column widths
        criteria_width = int(Inches(2.5))
        option_width = int(Inches(6.5 / len(options)))
        table.columns[0].width = criteria_width
        for col_idx in range(1, num_cols):
            table.columns[col_idx].width = option_width

        # Header row
        table.cell(0, 0).text = "Criteria"
        for col_idx, option in enumerate(options):
            cell = table.cell(0, col_idx + 1)
            cell.text = option
            cell.fill.solid()
            if col_idx + 1 == highlight_idx:
                cell.fill.fore_color.rgb = highlight_color
            else:
                cell.fill.fore_color.rgb = header_color
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        # Style criteria header
        cell = table.cell(0, 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Data rows
        for row_idx, criterion in enumerate(criteria):
            # Criteria name
            cell = table.cell(row_idx + 1, 0)
            cell.text = criterion.get("name", "")
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(10)

            # Values
            values = criterion.get("values", [])
            for col_idx, value in enumerate(values):
                if col_idx < len(options):
                    cell = table.cell(row_idx + 1, col_idx + 1)
                    cell.text = str(value)
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(10)
                    p.alignment = PP_ALIGN.CENTER

                    # Highlight column
                    if col_idx + 1 == highlight_idx:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 248, 225)

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added comparison matrix: {title} ({len(options)} options, {len(criteria)} criteria)")

    def add_timeline_slide(
        self,
        title: str,
        events: List[Dict[str, str]],
        notes: str = ""
    ) -> None:
        """
        Add a timeline slide with chronological events.

        Args:
            title: Slide title
            events: List of event dicts with 'date', 'label', and optional 'description'
            notes: Speaker notes
        """
        layout = self._get_layout(SlideType.TITLE_ONLY)
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        num_events = min(len(events), 6)
        events = events[:num_events]

        # Draw timeline line
        line_top = 3.0
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5),
            Inches(line_top),
            Inches(9.0),
            Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor.from_string(self.theme_colors.primary[1:])
        line.line.fill.background()

        # Add events
        event_width = 9.0 / num_events
        start_left = 0.5

        for i, event in enumerate(events):
            left = start_left + (i * event_width) + (event_width / 2) - 0.2

            # Event marker
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left),
                Inches(line_top - 0.15),
                Inches(0.35),
                Inches(0.35)
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = RGBColor.from_string(self.theme_colors.primary[1:])
            marker.line.fill.background()

            # Date above line
            date_box = slide.shapes.add_textbox(
                Inches(start_left + (i * event_width)),
                Inches(2.2),
                Inches(event_width - 0.1),
                Inches(0.6)
            )
            tf = date_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = event.get("date", "")
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor.from_string(self.theme_colors.primary[1:])
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Label below line
            label_box = slide.shapes.add_textbox(
                Inches(start_left + (i * event_width)),
                Inches(3.5),
                Inches(event_width - 0.1),
                Inches(0.5)
            )
            tf = label_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = event.get("label", "")
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Description
            if event.get("description"):
                desc_box = slide.shapes.add_textbox(
                    Inches(start_left + (i * event_width)),
                    Inches(4.0),
                    Inches(event_width - 0.1),
                    Inches(1.0)
                )
                tf = desc_box.text_frame
                tf.word_wrap = True
                tf.paragraphs[0].text = event["description"]
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added timeline slide: {title} ({num_events} events)")

    def add_key_value_slide(
        self,
        title: str,
        pairs: List[Dict[str, str]],
        notes: str = ""
    ) -> None:
        """
        Add a key-value pairs slide for statistics or facts.

        Args:
            title: Slide title
            pairs: List of dicts with 'key' and 'value'
            notes: Speaker notes
        """
        layout = self._get_layout(SlideType.TITLE_ONLY)
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        num_pairs = min(len(pairs), 6)
        pairs = pairs[:num_pairs]

        # Arrange in 2 or 3 columns
        cols = 3 if num_pairs > 4 else 2
        rows = (num_pairs + cols - 1) // cols

        card_width = 2.4 if cols == 3 else 3.2
        card_height = 1.4 if rows > 2 else 1.8
        h_gap = 0.2
        v_gap = 0.2
        start_left = (10.0 - (cols * card_width + (cols - 1) * h_gap)) / 2
        start_top = 1.8

        primary_color = RGBColor.from_string(self.theme_colors.primary[1:])

        for i, pair in enumerate(pairs):
            col = i % cols
            row = i // cols

            left = start_left + col * (card_width + h_gap)
            top = start_top + row * (card_height + v_gap)

            # Card background - use theme colors
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left),
                Inches(top),
                Inches(card_width),
                Inches(card_height)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self._hex_to_rgb(self.theme_colors.surface)
            card.line.color.rgb = self._hex_to_rgb(self._darken_color(self.theme_colors.surface, 0.15))
            card.line.width = Pt(1)

            # Value (large, prominent)
            value_box = slide.shapes.add_textbox(
                Inches(left + 0.1),
                Inches(top + 0.2),
                Inches(card_width - 0.2),
                Inches(card_height * 0.5)
            )
            tf = value_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = pair.get("value", "")
            tf.paragraphs[0].font.size = Pt(20)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = primary_color
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Key (label) - use secondary text color
            key_box = slide.shapes.add_textbox(
                Inches(left + 0.1),
                Inches(top + card_height * 0.55),
                Inches(card_width - 0.2),
                Inches(card_height * 0.4)
            )
            tf = key_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = pair.get("key", "")
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme_colors.text_secondary)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added key-value slide: {title} ({num_pairs} pairs)")

    def add_callout_slide(
        self,
        title: str,
        callouts: List[Dict[str, str]],
        notes: str = ""
    ) -> None:
        """
        Add a slide with stacked callout boxes.

        Args:
            title: Slide title
            callouts: List of callout dicts with 'callout_type', 'heading', 'text'
            notes: Speaker notes
        """
        layout = self._get_layout(SlideType.TITLE_ONLY)
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        num_callouts = min(len(callouts), 4)  # Max 4 callouts
        callouts = callouts[:num_callouts]

        callout_height = 1.0
        v_gap = 0.15
        start_top = 1.8
        start_left = 0.5
        width = 9.0

        for i, callout in enumerate(callouts):
            top = start_top + i * (callout_height + v_gap)
            callout_type = callout.get("callout_type", "info")
            heading = callout.get("heading", "")
            text = callout.get("text", "")

            # Get colors from theme-derived container colors
            colors = self._container_colors.get(callout_type, self._container_colors["info"])
            bg_color = colors["bg"]
            text_color = colors["text"]

            # Create callout box
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(start_left),
                Inches(top),
                Inches(width),
                Inches(callout_height)
            )

            # Set fill color
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._hex_to_rgb(bg_color)

            # Remove line
            shape.line.fill.background()

            # Add text
            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.3)
            tf.margin_right = Inches(0.2)
            tf.margin_top = Inches(0.15)
            tf.margin_bottom = Inches(0.15)

            # Add heading if provided
            if heading:
                p = tf.paragraphs[0]
                p.text = heading
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = self._hex_to_rgb(text_color)
                p.alignment = PP_ALIGN.LEFT

                # Add text as new paragraph
                if text:
                    p_text = tf.add_paragraph()
                    p_text.text = text
                    p_text.font.size = Pt(12)
                    p_text.font.color.rgb = self._hex_to_rgb(text_color)
                    p_text.alignment = PP_ALIGN.LEFT
                    p_text.space_before = Pt(6)
            else:
                # Just text, no heading
                tf.paragraphs[0].text = text
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].font.color.rgb = RGBColor.from_string(text_color[1:])
                tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added callout slide: {title} ({num_callouts} callouts)")

    def add_stats_grid_slide(
        self,
        title: str,
        stats: List[Dict[str, Any]],
        layout: str = "auto",
        notes: str = ""
    ) -> None:
        """
        Add a stats grid slide with value/label cards and optional trends.

        Args:
            title: Slide title
            stats: List of stat dicts with 'value', 'label', optional 'trend', 'highlight'
            layout: 'auto', '2col', '3col', or '4col'
            notes: Speaker notes
        """
        slide_layout = self._get_layout(SlideType.TITLE_ONLY)
        slide = self.prs.slides.add_slide(slide_layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        num_stats = min(len(stats), 8)  # Max 8 stats
        stats = stats[:num_stats]

        # Determine grid layout
        if layout == "auto":
            if num_stats <= 2:
                cols = 2
            elif num_stats <= 4:
                cols = 2
            elif num_stats <= 6:
                cols = 3
            else:
                cols = 4
        elif layout == "2col":
            cols = 2
        elif layout == "3col":
            cols = 3
        elif layout == "4col":
            cols = 4
        else:
            cols = 3  # Default

        rows = (num_stats + cols - 1) // cols

        # Calculate dimensions
        if cols == 2:
            card_width = 4.0
        elif cols == 3:
            card_width = 2.8
        else:
            card_width = 2.2

        card_height = 1.8 if rows > 2 else 2.2
        h_gap = 0.2
        v_gap = 0.2
        start_left = (10.0 - (cols * card_width + (cols - 1) * h_gap)) / 2
        start_top = 1.8

        primary_color = RGBColor.from_string(self.theme_colors.primary[1:])
        highlight_color = RGBColor.from_string(self.theme_colors.accent[1:])

        for i, stat in enumerate(stats):
            col = i % cols
            row = i // cols

            left = start_left + col * (card_width + h_gap)
            top = start_top + row * (card_height + v_gap)

            value = stat.get("value", "")
            label = stat.get("label", "")
            trend = stat.get("trend", None)  # "up", "down", or None
            highlight = stat.get("highlight", False)

            # Card background - use theme colors
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left),
                Inches(top),
                Inches(card_width),
                Inches(card_height)
            )
            card.fill.solid()
            if highlight:
                card.fill.fore_color.rgb = self._hex_to_rgb(self._lighten_color(self.theme_colors.accent, 0.85))
            else:
                card.fill.fore_color.rgb = self._hex_to_rgb(self.theme_colors.surface)
            card.line.color.rgb = self._hex_to_rgb(self._darken_color(self.theme_colors.surface, 0.15))
            card.line.width = Pt(1)

            # Value (large, prominent - 36pt)
            value_box = slide.shapes.add_textbox(
                Inches(left + 0.1),
                Inches(top + 0.2),
                Inches(card_width - 0.2),
                Inches(card_height * 0.5)
            )
            tf = value_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = str(value)
            tf.paragraphs[0].font.size = Pt(36)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = highlight_color if highlight else primary_color
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Trend indicator (if provided) - use semantic colors
            if trend:
                trend_symbol = "▲" if trend == "up" else "▼"
                trend_color = self._hex_to_rgb(self.theme_colors.success) if trend == "up" else self._hex_to_rgb(self.theme_colors.error)

                trend_box = slide.shapes.add_textbox(
                    Inches(left + card_width - 0.5),
                    Inches(top + 0.15),
                    Inches(0.4),
                    Inches(0.4)
                )
                tf_trend = trend_box.text_frame
                tf_trend.paragraphs[0].text = trend_symbol
                tf_trend.paragraphs[0].font.size = Pt(20)
                tf_trend.paragraphs[0].font.color.rgb = trend_color
                tf_trend.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Label - use secondary text color
            label_box = slide.shapes.add_textbox(
                Inches(left + 0.1),
                Inches(top + card_height * 0.55),
                Inches(card_width - 0.2),
                Inches(card_height * 0.4)
            )
            tf = label_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = label
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme_colors.text_secondary)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added stats grid slide: {title} ({num_stats} stats, {cols} columns)")
    def add_data_viz_slide(self, slide_data: Dict[str, Any]) -> None:
        """
        Add a data visualization slide with charts.

        Args:
            slide_data: Slide data dict containing title, content with chart configuration

        Expected content structure:
        {
            "type": "data_viz",
            "title": "Chart Title",
            "content": {
                "chart_type": "column_clustered|column_stacked|bar_clustered|bar_stacked|line|line_markers|pie|area|area_stacked",
                "chart_data": {
                    "categories": ["Q1", "Q2", "Q3", "Q4"],
                    "series": [
                        {"name": "Series 1", "values": [120, 145, 160, 185]},
                        {"name": "Series 2", "values": [135, 168, 195, 220]}
                    ]
                },
                "chart_title": "Optional chart title",
                "show_legend": true,
                "legend_position": "right|bottom|top|left",
                "show_data_labels": false,
                "alt_text": "Alternative text for accessibility"
            }
        }
        """
        title = slide_data.get("title", "Data Visualization")
        content = slide_data.get("content", {})
        notes = slide_data.get("notes", "")

        # Get chart configuration
        chart_type_str = content.get("chart_type", "column_clustered")
        chart_data_raw = content.get("chart_data", {})
        chart_title = content.get("chart_title", "")
        show_legend = content.get("show_legend", True)
        legend_position_str = content.get("legend_position", "right")
        show_data_labels = content.get("show_data_labels", False)
        alt_text = content.get("alt_text", "")

        # Map chart type string to XL_CHART_TYPE enum
        chart_type_map = {
            "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
            "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
            "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
            "line": XL_CHART_TYPE.LINE,
            "line_markers": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
            "area": XL_CHART_TYPE.AREA,
            "area_stacked": XL_CHART_TYPE.AREA_STACKED,
        }

        # Map legend position string to XL_LEGEND_POSITION enum
        legend_position_map = {
            "right": XL_LEGEND_POSITION.RIGHT,
            "bottom": XL_LEGEND_POSITION.BOTTOM,
            "top": XL_LEGEND_POSITION.TOP,
            "left": XL_LEGEND_POSITION.LEFT,
        }

        # Validate chart type
        chart_type = chart_type_map.get(chart_type_str)
        if not chart_type:
            logger.error(f"Invalid chart type: {chart_type_str}. Defaulting to column_clustered.")
            chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED

        # Create slide with Title Only layout
        layout = self._get_layout(SlideType.TITLE_ONLY)
        slide = self.prs.slides.add_slide(layout)

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = title

        try:
            # Validate chart data
            categories = chart_data_raw.get("categories", [])
            series_list = chart_data_raw.get("series", [])

            if not categories:
                logger.error("No categories provided for chart")
                categories = ["Category 1"]

            if not series_list:
                logger.error("No series data provided for chart")
                series_list = [{"name": "Series 1", "values": [0] * len(categories)}]

            # Create CategoryChartData object
            chart_data = CategoryChartData()
            chart_data.categories = categories

            # Add series data
            # For pie charts, only use the first series
            if chart_type == XL_CHART_TYPE.PIE:
                if series_list:
                    first_series = series_list[0]
                    series_name = first_series.get("name", "Data")
                    series_values = first_series.get("values", [])

                    # Ensure values match categories length
                    if len(series_values) < len(categories):
                        series_values.extend([0] * (len(categories) - len(series_values)))
                    elif len(series_values) > len(categories):
                        series_values = series_values[:len(categories)]

                    chart_data.add_series(series_name, series_values)
            else:
                # For non-pie charts, add all series
                for series in series_list:
                    series_name = series.get("name", "Series")
                    series_values = series.get("values", [])

                    # Ensure values match categories length
                    if len(series_values) < len(categories):
                        series_values.extend([0] * (len(categories) - len(series_values)))
                    elif len(series_values) > len(categories):
                        series_values = series_values[:len(categories)]

                    chart_data.add_series(series_name, series_values)

            # Add chart to slide
            # Position: x=0.5", y=1.8", width=9.0", height=5.0"
            x = Inches(0.5)
            y = Inches(1.8)
            width = Inches(9.0)
            height = Inches(5.0)

            graphic_frame = slide.shapes.add_chart(
                chart_type, x, y, width, height, chart_data
            )

            chart = graphic_frame.chart

            # Set chart title if provided
            if chart_title:
                chart.has_title = True
                chart.chart_title.text_frame.text = chart_title

            # Configure legend
            if show_legend and chart.has_legend:
                legend_position = legend_position_map.get(legend_position_str, XL_LEGEND_POSITION.RIGHT)
                chart.legend.position = legend_position

                # Position legend based on number of series for better layout
                num_series = len(series_list)
                if num_series > 4 and legend_position_str == "right":
                    # Keep right position for many series
                    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            elif not show_legend and chart.has_legend:
                chart.has_legend = False

            # Configure data labels
            if show_data_labels:
                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels
                data_labels.font.size = Pt(10)

            # Set alt text for accessibility
            if alt_text:
                graphic_frame.name = alt_text

        except Exception as e:
            logger.error(f"Error creating chart: {e}")
            # Add error message to slide as fallback
            error_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(3.0), Inches(8.0), Inches(1.0)
            )
            tf = error_box.text_frame
            tf.text = f"Error creating chart: {str(e)}"
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.color.rgb = RGBColor(220, 53, 69)  # Red

        # Add speaker notes
        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added data visualization slide: {title} ({chart_type_str})")




    # ==================== VISUAL STYLING METHODS ====================

    def add_callout_box(
        self,
        slide,
        text: str,
        callout_type: str = "info",
        left: float = 0.5,
        top: float = 5.0,
        width: float = 9.0,
        height: float = 1.0
    ) -> None:
        """
        Add a callout box to a slide.

        Args:
            slide: The slide object
            text: Text content for the callout
            callout_type: info (blue), success (green), warning (yellow), tip (purple)
            left, top, width, height: Position in inches
        """
        colors = {
            "info": ("#e3f2fd", "#1565c0"),      # Light blue bg, dark blue text
            "success": ("#e8f5e9", "#2e7d32"),   # Light green bg, dark green text
            "warning": ("#fff8e1", "#f57f17"),   # Light yellow bg, dark yellow text
            "tip": ("#f3e5f5", "#7b1fa2"),       # Light purple bg, dark purple text
        }
        bg_color, text_color = colors.get(callout_type, colors["info"])

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )

        # Set fill color
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(bg_color[1:])

        # Remove line
        shape.line.fill.background()

        # Add text
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = RGBColor.from_string(text_color[1:])
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        # Add left margin for icon space
        tf.margin_left = Inches(0.3)

    def add_rounded_card(
        self,
        slide,
        title: str,
        content_items: List[str],
        left: float = 0.5,
        top: float = 1.5,
        width: float = 4.0,
        height: float = 4.0
    ) -> None:
        """
        Add a rounded card container with title and bullet content.

        Args:
            slide: The slide object
            title: Card title
            content_items: List of bullet points
            left, top, width, height: Position in inches
        """
        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(248, 249, 250)  # Light gray
        card.line.color.rgb = RGBColor(222, 226, 230)  # Border gray
        card.line.width = Pt(1)

        # Card title
        title_box = slide.shapes.add_textbox(
            Inches(left + 0.2), Inches(top + 0.2),
            Inches(width - 0.4), Inches(0.5)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)

        # Card content
        content_box = slide.shapes.add_textbox(
            Inches(left + 0.2), Inches(top + 0.7),
            Inches(width - 0.4), Inches(height - 0.9)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(content_items):
            if i == 0:
                tf.paragraphs[0].text = f"• {item}"
            else:
                p = tf.add_paragraph()
                p.text = f"• {item}"
            tf.paragraphs[i].font.size = Pt(12)
            tf.paragraphs[i].font.color.rgb = RGBColor(51, 51, 51)

    def add_speech_bubble(
        self,
        slide,
        text: str,
        attribution: str = "",
        left: float = 1.0,
        top: float = 2.0,
        width: float = 8.0,
        height: float = 2.5
    ) -> None:
        """
        Add a speech bubble for quotes.

        Args:
            slide: The slide object
            text: Quote text
            attribution: Quote attribution
            left, top, width, height: Position in inches
        """
        # Bubble shape
        bubble = slide.shapes.add_shape(
            MSO_SHAPE.OVAL_CALLOUT,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = RGBColor(227, 242, 253)  # Light blue
        bubble.line.color.rgb = RGBColor(33, 150, 243)  # Blue border
        bubble.line.width = Pt(2)

        # Quote text
        tf = bubble.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = f'"{text}"'
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.italic = True
        tf.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Attribution below bubble
        if attribution:
            attr_box = slide.shapes.add_textbox(
                Inches(left), Inches(top + height + 0.2),
                Inches(width), Inches(0.4)
            )
            attr_tf = attr_box.text_frame
            attr_tf.paragraphs[0].text = f"— {attribution}"
            attr_tf.paragraphs[0].font.size = Pt(14)
            attr_tf.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
            attr_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def add_circle_badge(
        self,
        slide,
        text: str,
        left: float = 0.5,
        top: float = 1.5,
        size: float = 0.6,
        color: str = "#2c5aa0"
    ) -> None:
        """
        Add a circular badge with number or short text.

        Args:
            slide: The slide object
            text: Badge text (number or short label)
            left, top: Position in inches
            size: Diameter in inches
            color: Hex color for badge
        """
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left), Inches(top), Inches(size), Inches(size)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor.from_string(color[1:])
        circle.line.fill.background()

        # Center text in circle
        tf = circle.text_frame
        tf.paragraphs[0].text = str(text)
        tf.paragraphs[0].font.size = Pt(int(size * 24))
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

    def add_styled_content_slide(
        self,
        title: str,
        bullets: List[str],
        step_number: Optional[int] = None,
        callout_text: Optional[str] = None,
        callout_type: str = "tip",
        notes: str = ""
    ) -> None:
        """
        Add a styled content slide with visual elements.

        Args:
            title: Slide title
            bullets: List of bullet points
            step_number: Optional step number for circular badge
            callout_text: Optional callout box text
            callout_type: Type of callout (info, success, warning, tip)
            notes: Speaker notes
        """
        layout = self._get_layout(SlideType.CONTENT)
        slide = self.prs.slides.add_slide(layout)

        # Clear default content
        if slide.shapes.title:
            slide.shapes.title.text = title

        # Add step number badge if provided
        if step_number is not None:
            self.add_circle_badge(slide, str(step_number), left=0.3, top=0.3, size=0.5)

        # Add bullets in a rounded card
        self.add_rounded_card(
            slide, "Key Points", bullets,
            left=0.5, top=1.8, width=5.5, height=3.5
        )

        # Add callout box if provided
        if callout_text:
            self.add_callout_box(
                slide, callout_text, callout_type,
                left=6.2, top=1.8, width=3.3, height=1.5
            )

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added styled content slide: {title}")

    def add_cards_grid_slide(
        self,
        title: str,
        cards: List[Dict[str, Any]],
        columns: int = 2,
        notes: str = ""
    ) -> None:
        """
        Add a cards grid slide with multiple content cards.

        Args:
            title: Slide title
            cards: List of card dicts with 'title', 'items' (list), and optional 'color'
            columns: Number of columns (2 or 3)
            notes: Speaker notes
        """
        layout = self._get_layout(SlideType.TITLE_ONLY)
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        num_cards = min(len(cards), 6)
        cards = cards[:num_cards]

        # Validate columns
        columns = 3 if columns == 3 else 2

        # Calculate layout
        rows = (num_cards + columns - 1) // columns
        card_width = 2.4 if columns == 3 else 3.2
        card_height = 2.5 if rows > 1 else 3.2
        h_gap = 0.3
        v_gap = 0.3
        start_left = (10.0 - (columns * card_width + (columns - 1) * h_gap)) / 2
        start_top = 1.8

        # Color mapping
        color_map = {
            "default": RGBColor(248, 249, 250),
            "primary": RGBColor.from_string(self.theme_colors.primary[1:]),
            "secondary": RGBColor.from_string(self.theme_colors.secondary[1:]),
            "accent": RGBColor.from_string(self.theme_colors.accent[1:]),
        }

        for i, card_data in enumerate(cards):
            col = i % columns
            row = i // columns

            left = start_left + col * (card_width + h_gap)
            top = start_top + row * (card_height + v_gap)

            # Get card color
            card_color_key = card_data.get("color", "default")
            card_bg_color = color_map.get(card_color_key, color_map["default"])

            # Determine text color based on background
            if card_color_key in ["primary", "secondary", "accent"]:
                text_color = RGBColor(255, 255, 255)
                title_color = RGBColor(255, 255, 255)
            else:
                text_color = RGBColor(51, 51, 51)
                title_color = RGBColor.from_string(self.theme_colors.primary[1:])

            card_title = card_data.get("title", "")
            card_items = card_data.get("items", [])

            # Card background
            card_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left),
                Inches(top),
                Inches(card_width),
                Inches(card_height)
            )
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = card_bg_color
            card_shape.line.color.rgb = RGBColor(222, 226, 230)
            card_shape.line.width = Pt(1)

            # Card title
            title_box = slide.shapes.add_textbox(
                Inches(left + 0.2),
                Inches(top + 0.2),
                Inches(card_width - 0.4),
                Inches(0.5)
            )
            tf = title_box.text_frame
            tf.paragraphs[0].text = card_title
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = title_color
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT

            # Card content
            content_box = slide.shapes.add_textbox(
                Inches(left + 0.2),
                Inches(top + 0.8),
                Inches(card_width - 0.4),
                Inches(card_height - 1.0)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            for j, item in enumerate(card_items[:6]):  # Max 6 items per card
                if j == 0:
                    tf.paragraphs[0].text = f"• {item}"
                    tf.paragraphs[0].font.size = Pt(16)
                    tf.paragraphs[0].font.color.rgb = text_color
                else:
                    p = tf.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(16)
                    p.font.color.rgb = text_color

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added cards grid slide: {title} ({num_cards} cards, {columns} columns)")

    def add_key_concept_slide(
        self,
        title: str,
        concept_title: str,
        definition: str,
        details: Optional[List[str]] = None,
        concept_style: str = "boxed",
        notes: str = ""
    ) -> None:
        """
        Add a key concept slide with prominent concept display.

        Args:
            title: Slide title
            concept_title: The concept name (displayed prominently)
            definition: Concept definition
            details: Optional list of additional details/bullets
            concept_style: "boxed" (default), "highlighted", or "minimal"
            notes: Speaker notes
        """
        layout = self._get_layout(SlideType.TITLE_ONLY)
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        primary_color = RGBColor.from_string(self.theme_colors.primary[1:])

        # Concept box positioning
        box_left = 1.0
        box_top = 2.0
        box_width = 8.0
        box_height = 2.0

        if concept_style == "boxed":
            # Prominent colored box with white text
            concept_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(box_left),
                Inches(box_top),
                Inches(box_width),
                Inches(box_height)
            )
            concept_box.fill.solid()
            concept_box.fill.fore_color.rgb = primary_color
            concept_box.line.fill.background()

            # Concept title (large, bold, white)
            title_box = slide.shapes.add_textbox(
                Inches(box_left + 0.3),
                Inches(box_top + 0.3),
                Inches(box_width - 0.6),
                Inches(0.6)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = concept_title
            tf.paragraphs[0].font.size = Pt(24)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT

            # Definition (white text below title)
            def_box = slide.shapes.add_textbox(
                Inches(box_left + 0.3),
                Inches(box_top + 1.0),
                Inches(box_width - 0.6),
                Inches(box_height - 1.3)
            )
            tf = def_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = definition
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        elif concept_style == "highlighted":
            # Highlighted concept with colored background, darker text
            highlight_color = RGBColor(227, 242, 253)  # Light blue

            concept_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(box_left),
                Inches(box_top),
                Inches(box_width),
                Inches(box_height)
            )
            concept_box.fill.solid()
            concept_box.fill.fore_color.rgb = highlight_color
            concept_box.line.color.rgb = primary_color
            concept_box.line.width = Pt(2)

            # Concept title (colored)
            title_box = slide.shapes.add_textbox(
                Inches(box_left + 0.3),
                Inches(box_top + 0.3),
                Inches(box_width - 0.6),
                Inches(0.6)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = concept_title
            tf.paragraphs[0].font.size = Pt(24)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = primary_color
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT

            # Definition (dark text)
            def_box = slide.shapes.add_textbox(
                Inches(box_left + 0.3),
                Inches(box_top + 1.0),
                Inches(box_width - 0.6),
                Inches(box_height - 1.3)
            )
            tf = def_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = definition
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        else:  # minimal
            # Simple text-based concept display
            # Concept title (large, colored, no box)
            title_box = slide.shapes.add_textbox(
                Inches(box_left),
                Inches(box_top),
                Inches(box_width),
                Inches(0.6)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = concept_title
            tf.paragraphs[0].font.size = Pt(28)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = primary_color
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT

            # Definition (normal text)
            def_box = slide.shapes.add_textbox(
                Inches(box_left),
                Inches(box_top + 0.8),
                Inches(box_width),
                Inches(1.2)
            )
            tf = def_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = definition
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        # Add details as bullets below the concept box
        if details:
            details_top = box_top + box_height + 0.5
            details_box = slide.shapes.add_textbox(
                Inches(box_left),
                Inches(details_top),
                Inches(box_width),
                Inches(5.5 - details_top)
            )
            tf = details_box.text_frame
            tf.word_wrap = True

            for i, detail in enumerate(details[:5]):  # Max 5 details
                if i == 0:
                    tf.paragraphs[0].text = f"• {detail}"
                    tf.paragraphs[0].font.size = Pt(13)
                    tf.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
                else:
                    p = tf.add_paragraph()
                    p.text = f"• {detail}"
                    p.font.size = Pt(13)
                    p.font.color.rgb = RGBColor(51, 51, 51)

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added key concept slide: {title} ({concept_style} style)")

    def add_image_grid_slide(self, slide_data: Dict[str, Any]) -> None:
        """
        Add an image grid slide with multiple images arranged in a grid layout.

        Args:
            slide_data: Dict with 'type', 'title', 'content', 'notes'
                content should contain:
                    - images: List of dicts with 'path', 'caption', 'alt_text'
                    - grid_layout: 'auto', '2x1', '1x2', '2x2', '3x1', '1x3', '3x2', '2x3', '3x3'
                    - show_captions: bool (default True)
                    - gap_size: 'small', 'medium', 'large' (default 'medium')
        """
        title = slide_data.get("title", "")
        content = slide_data.get("content", {})
        notes = slide_data.get("notes", "")

        images = content.get("images", [])
        grid_layout = content.get("grid_layout", "auto")
        show_captions = content.get("show_captions", True)
        gap_size = content.get("gap_size", "medium")

        # Validate image count
        num_images = len(images)
        if num_images < 2:
            logger.warning(f"Image grid requires at least 2 images, got {num_images}")
            return
        if num_images > 9:
            logger.warning(f"Image grid supports max 9 images, truncating from {num_images}")
            images = images[:9]
            num_images = 9

        # Create slide with Title Only layout
        layout = self._get_layout(SlideType.TITLE_ONLY)
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Determine grid dimensions (rows x cols)
        if grid_layout == "auto":
            # Auto-detect based on image count
            if num_images == 2:
                rows, cols = 1, 2
            elif num_images == 3:
                rows, cols = 1, 3
            elif num_images == 4:
                rows, cols = 2, 2
            elif num_images <= 6:
                rows, cols = 2, 3
            else:  # 7-9 images
                rows, cols = 3, 3
        else:
            # Parse explicit layout (e.g., "2x3" -> rows=2, cols=3)
            try:
                parts = grid_layout.lower().split('x')
                if len(parts) == 2:
                    rows, cols = int(parts[0]), int(parts[1])
                else:
                    logger.warning(f"Invalid grid_layout '{grid_layout}', using auto")
                    rows, cols = 2, 2
            except (ValueError, AttributeError):
                logger.warning(f"Invalid grid_layout '{grid_layout}', using auto")
                rows, cols = 2, 2

        # Validate grid can hold all images
        if rows * cols < num_images:
            logger.warning(f"Grid {rows}x{cols} cannot hold {num_images} images, using auto layout")
            # Fall back to auto
            if num_images <= 4:
                rows, cols = 2, 2
            elif num_images <= 6:
                rows, cols = 2, 3
            else:
                rows, cols = 3, 3

        # Gap size mapping (in inches)
        gap_sizes = {
            "small": 0.1,
            "medium": 0.2,
            "large": 0.3
        }
        gap = gap_sizes.get(gap_size, 0.2)

        # Grid area dimensions (in inches)
        grid_left = 0.5
        grid_top = 1.8
        grid_width = 9.0
        grid_height = 5.0

        # Caption height (if enabled)
        caption_height = 0.3 if show_captions else 0.0

        # Calculate cell dimensions
        cell_width = (grid_width - gap * (cols - 1)) / cols
        cell_height = (grid_height - gap * (rows - 1)) / rows

        # Image height (leaving space for caption)
        image_height = cell_height - caption_height

        # Place images in grid
        for i, image_data in enumerate(images):
            if i >= rows * cols:
                break  # Grid is full

            row = i // cols
            col = i % cols

            # Calculate position
            left = grid_left + col * (cell_width + gap)
            top = grid_top + row * (cell_height + gap)

            image_path = image_data.get("path", "")
            caption = image_data.get("caption", "")
            alt_text = image_data.get("alt_text", "")

            # Validate image file exists
            if not image_path:
                logger.warning(f"Image {i+1} has no path, skipping")
                continue

            image_file = Path(image_path)
            if not image_file.exists():
                logger.warning(f"Image not found: {image_path}, skipping")
                continue

            try:
                # Add image
                pic = slide.shapes.add_picture(
                    str(image_file),
                    Inches(left),
                    Inches(top),
                    width=Inches(cell_width),
                    height=Inches(image_height)
                )

                # Set alt text for accessibility
                if alt_text:
                    pic.name = alt_text
                elif caption:
                    pic.name = caption
                else:
                    pic.name = f"Image {i+1}"

                # Add caption if enabled
                if show_captions and caption:
                    caption_top = top + image_height
                    caption_box = slide.shapes.add_textbox(
                        Inches(left),
                        Inches(caption_top),
                        Inches(cell_width),
                        Inches(caption_height)
                    )
                    tf = caption_box.text_frame
                    tf.word_wrap = True
                    tf.paragraphs[0].text = caption
                    tf.paragraphs[0].font.size = Pt(10)
                    tf.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            except Exception as e:
                logger.warning(f"Failed to add image {image_path}: {e}")
                continue

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added image grid slide: {title} ({num_images} images in {rows}x{cols} grid)")

    def add_agenda_slide(self, slide_data: Dict[str, Any]) -> None:
        """
        Add an agenda slide with navigation and progress tracking.

        Args:
            slide_data: Slide data dictionary with type, title, content, and notes

        Content structure:
            {
                "agenda_items": [
                    {
                        "section": "Section name",
                        "description": "Optional time/details",
                        "status": "complete" | "current" | "upcoming"
                    }
                ],
                "agenda_style": "numbered" | "checkmarks" | "progress_bar",
                "show_progress": true/false
            }
        """
        title = slide_data.get("title", "Agenda")
        content = slide_data.get("content", {})
        notes = slide_data.get("notes", "")

        layout = self._get_layout(SlideType.TITLE_ONLY)
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Get agenda items and settings
        agenda_items = content.get("agenda_items", [])
        if not agenda_items:
            logger.warning("No agenda items provided for agenda slide")
            return

        # Limit to 2-8 items
        num_items = min(max(len(agenda_items), 2), 8)
        agenda_items = agenda_items[:num_items]

        agenda_style = content.get("agenda_style", "numbered")
        show_progress = content.get("show_progress", False)

        # Colors
        primary_color = RGBColor.from_string(self.theme_colors.primary[1:])
        success_color = RGBColor(46, 125, 50)  # Green
        current_color = RGBColor(255, 193, 7)  # Yellow/amber
        upcoming_color = RGBColor(158, 158, 158)  # Gray
        text_dark = RGBColor(51, 51, 51)

        # Layout settings
        start_left = 1.5
        start_top = 2.2
        item_height = 0.6
        item_spacing = 0.15
        indicator_width = 0.4

        # Calculate total height and adjust spacing if needed
        total_height = num_items * (item_height + item_spacing)
        max_height = 4.5 if show_progress else 5.0
        if total_height > max_height:
            item_spacing = (max_height - (num_items * item_height)) / (num_items - 1)

        # Track current item for highlighting
        current_index = -1
        for i, item in enumerate(agenda_items):
            if item.get("status") == "current":
                current_index = i
                break

        # Add agenda items
        for i, item in enumerate(agenda_items):
            top = start_top + i * (item_height + item_spacing)
            section_name = item.get("section", f"Item {i + 1}")
            description = item.get("description", "")
            status = item.get("status", "upcoming")

            # Determine indicator and color based on status
            if status == "complete":
                indicator = "✓"
                indicator_color = success_color
            elif status == "current":
                indicator = "→"
                indicator_color = current_color
            else:  # upcoming
                indicator = "○"
                indicator_color = upcoming_color

            # Override indicator based on style
            if agenda_style == "numbered":
                indicator = str(i + 1)
                if status == "complete":
                    indicator_color = success_color
                elif status == "current":
                    indicator_color = primary_color
                else:
                    indicator_color = upcoming_color
            elif agenda_style == "progress_bar":
                if status == "complete":
                    indicator = "■"
                    indicator_color = success_color
                elif status == "current":
                    indicator = "▶"
                    indicator_color = primary_color
                else:
                    indicator = "□"
                    indicator_color = upcoming_color

            # Highlight background for current item
            if status == "current":
                highlight_box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(start_left - 0.2),
                    Inches(top - 0.05),
                    Inches(7.0),
                    Inches(item_height + 0.1)
                )
                highlight_box.fill.solid()
                highlight_box.fill.fore_color.rgb = RGBColor(255, 248, 225)  # Light yellow
                highlight_box.line.fill.background()

            # Indicator/number box
            indicator_box = slide.shapes.add_textbox(
                Inches(start_left),
                Inches(top),
                Inches(indicator_width),
                Inches(item_height)
            )
            tf = indicator_box.text_frame
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.paragraphs[0].text = indicator
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = indicator_color
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Section name
            section_left = start_left + indicator_width + 0.2
            section_width = 4.5 if description else 6.0

            section_box = slide.shapes.add_textbox(
                Inches(section_left),
                Inches(top),
                Inches(section_width),
                Inches(item_height)
            )
            tf = section_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.paragraphs[0].text = section_name
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.bold = (status == "current")
            tf.paragraphs[0].font.color.rgb = text_dark if status == "current" else RGBColor(80, 80, 80)
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT

            # Description (optional - time, duration, etc.)
            if description:
                desc_left = section_left + section_width + 0.2
                desc_box = slide.shapes.add_textbox(
                    Inches(desc_left),
                    Inches(top),
                    Inches(1.5),
                    Inches(item_height)
                )
                tf = desc_box.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.paragraphs[0].text = description
                tf.paragraphs[0].font.size = Pt(11)
                tf.paragraphs[0].font.italic = True
                tf.paragraphs[0].font.color.rgb = RGBColor(120, 120, 120)
                tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Add progress bar at bottom if requested
        if show_progress:
            completed_count = sum(1 for item in agenda_items if item.get("status") == "complete")
            progress_pct = (completed_count / num_items) * 100

            progress_top = 6.8
            progress_left = 1.5
            progress_width = 7.0
            progress_height = 0.3

            # Background bar (gray)
            bg_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(progress_left),
                Inches(progress_top),
                Inches(progress_width),
                Inches(progress_height)
            )
            bg_bar.fill.solid()
            bg_bar.fill.fore_color.rgb = RGBColor(220, 220, 220)
            bg_bar.line.fill.background()

            # Progress bar (colored)
            if progress_pct > 0:
                progress_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(progress_left),
                    Inches(progress_top),
                    Inches(progress_width * (progress_pct / 100)),
                    Inches(progress_height)
                )
                progress_bar.fill.solid()
                progress_bar.fill.fore_color.rgb = success_color
                progress_bar.line.fill.background()

            # Progress percentage label
            pct_box = slide.shapes.add_textbox(
                Inches(progress_left + progress_width + 0.2),
                Inches(progress_top - 0.05),
                Inches(0.8),
                Inches(0.4)
            )
            tf = pct_box.text_frame
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.paragraphs[0].text = f"{int(progress_pct)}%"
            tf.paragraphs[0].font.size = Pt(11)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = primary_color
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added agenda slide: {title} ({num_items} items, {agenda_style} style)")

    def add_content_with_containers_slide(
        self,
        title: str,
        bullets: Optional[List[str]] = None,
        containers: Optional[List[Dict[str, Any]]] = None,
        notes: str = ""
    ) -> None:
        """
        Add a content slide with optional bullets and inline containers.

        Args:
            title: Slide title
            bullets: Optional list of bullet points
            containers: Optional list of container dicts with 'container_type' and data
            notes: Speaker notes
        """
        # Choose layout based on content
        if bullets:
            layout = self._get_layout(SlideType.CONTENT)
        else:
            layout = self._get_layout(SlideType.TITLE_ONLY)

        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Add bullets if present
        if bullets:
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Content placeholder
                    content_placeholder = shape
                    break

            if content_placeholder:
                tf = content_placeholder.text_frame
                tf.clear()

                for i, bullet in enumerate(bullets):
                    if i == 0:
                        tf.paragraphs[0].text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0

        # Position containers below bullets (or in main content area if no bullets)
        if containers:
            container_top = 4.5 if bullets else 2.0
            container_left = 0.5
            container_width = 9.0

            for i, container in enumerate(containers):
                container_type = container.get("container_type", "")

                if container_type == "callout":
                    self.add_callout_box(
                        slide,
                        text=container.get("text", ""),
                        callout_type=container.get("callout_type", "info"),
                        left=container_left,
                        top=container_top + (i * 1.2),
                        width=container_width,
                        height=1.0
                    )

                elif container_type == "stat_box":
                    # Single stat box inline
                    card_width = 3.0
                    card_height = 1.2
                    left = container_left + (container_width - card_width) / 2
                    top = container_top + (i * 1.4)

                    card = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(left),
                        Inches(top),
                        Inches(card_width),
                        Inches(card_height)
                    )
                    card.fill.solid()
                    card.fill.fore_color.rgb = RGBColor(248, 249, 250)
                    card.line.color.rgb = RGBColor(222, 226, 230)
                    card.line.width = Pt(1)

                    # Value
                    value_box = slide.shapes.add_textbox(
                        Inches(left + 0.1),
                        Inches(top + 0.1),
                        Inches(card_width - 0.2),
                        Inches(card_height * 0.5)
                    )
                    tf = value_box.text_frame
                    tf.paragraphs[0].text = container.get("value", "")
                    tf.paragraphs[0].font.size = Pt(20)
                    tf.paragraphs[0].font.bold = True
                    tf.paragraphs[0].font.color.rgb = RGBColor.from_string(self.theme_colors.primary[1:])
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

                    # Label
                    label_box = slide.shapes.add_textbox(
                        Inches(left + 0.1),
                        Inches(top + card_height * 0.5),
                        Inches(card_width - 0.2),
                        Inches(card_height * 0.45)
                    )
                    tf = label_box.text_frame
                    tf.paragraphs[0].text = container.get("label", "")
                    tf.paragraphs[0].font.size = Pt(10)
                    tf.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

                elif container_type == "card":
                    self.add_rounded_card(
                        slide,
                        title=container.get("title", ""),
                        content_items=container.get("items", []),
                        left=container_left,
                        top=container_top + (i * 1.5),
                        width=container_width,
                        height=1.3
                    )

                elif container_type == "concept_block":
                    # Render term + definition box
                    block_height = 1.2
                    top = container_top + (i * 1.4)

                    concept_box = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(container_left),
                        Inches(top),
                        Inches(container_width),
                        Inches(block_height)
                    )
                    concept_box.fill.solid()
                    concept_box.fill.fore_color.rgb = RGBColor(248, 249, 250)
                    concept_box.line.color.rgb = RGBColor(222, 226, 230)
                    concept_box.line.width = Pt(1)

                    tf = concept_box.text_frame
                    tf.word_wrap = True
                    tf.margin_left = Inches(0.2)
                    tf.margin_top = Inches(0.1)

                    # Term (bold)
                    tf.paragraphs[0].text = container.get("term", "")
                    tf.paragraphs[0].font.size = Pt(14)
                    tf.paragraphs[0].font.bold = True
                    tf.paragraphs[0].font.color.rgb = RGBColor.from_string(self.theme_colors.primary[1:])

                    # Definition
                    p = tf.add_paragraph()
                    p.text = container.get("definition", "")
                    p.font.size = Pt(11)
                    p.font.color.rgb = RGBColor(51, 51, 51)

        if notes:
            self._add_notes(slide, notes)

        self._slide_count += 1
        logger.info(f"Added content with containers slide: {title}")

    # ==================== END VISUAL STYLING METHODS ====================

    def create_from_structure(self, content: Dict[str, Any]) -> None:
        """
        Create presentation from structured content dictionary.

        Args:
            content: Structured content following the presentation schema

        Expected structure:
        {
            "metadata": {"title": "", "subtitle": "", "author": "", "date": ""},
            "sections": [
                {
                    "title": "Section Name",
                    "slides": [
                        {
                            "type": "title|content|comparison|...",
                            "title": "",
                            "content": {...}
                        }
                    ]
                }
            ]
        }
        """
        # Set metadata (title slide is now explicit in sections)
        if "metadata" in content:
            meta = content["metadata"]
            self.set_metadata(PresentationMetadata(
                title=meta.get("title", "Untitled"),
                subtitle=meta.get("subtitle", ""),
                author=meta.get("author", ""),
                date=meta.get("date", ""),
                subject=meta.get("subject", "")
            ))

        # Process sections
        for section in content.get("sections", []):
            # Process slides in section (section headers are explicit slides)
            for slide_data in section.get("slides", []):
                self._add_slide_from_data(slide_data)

        logger.info(f"Created presentation with {self._slide_count} slides")

    def _add_slide_from_data(self, slide_data: Dict[str, Any]) -> None:
        """Add a slide from structured data."""
        slide_type = slide_data.get("type", "content")
        title = slide_data.get("title", "")
        content = slide_data.get("content", {})
        notes = slide_data.get("notes", "")

        if slide_type == "title":
            self.add_title_slide(
                title=title,
                subtitle=content.get("subtitle", ""),
                notes=notes
            )

        elif slide_type == "section_header":
            self.add_section_header(
                title=title,
                subtitle=content.get("subtitle", ""),
                notes=notes
            )

        elif slide_type == "divider":
            self.add_divider_slide(slide_data)

        elif slide_type == "agenda":
            self.add_agenda_slide(slide_data)

        elif slide_type == "content":
            # Check if containers are present
            if content.get("containers"):
                self.add_content_with_containers_slide(
                    title=title,
                    bullets=content.get("bullets"),
                    containers=content.get("containers"),
                    notes=notes
                )
            else:
                self.add_content_slide(
                    title=title,
                    bullets=content.get("bullets", []),
                    notes=notes
                )

        elif slide_type == "two_content":
            self.add_two_content_slide(
                title=title,
                left_content=content.get("left", []),
                right_content=content.get("right", []),
                notes=notes
            )

        elif slide_type == "comparison":
            self.add_comparison_slide(
                title=title,
                left_title=content.get("left_title", ""),
                left_content=content.get("left", []),
                right_title=content.get("right_title", ""),
                right_content=content.get("right", []),
                notes=notes
            )

        elif slide_type == "image":
            self.add_image_slide(
                title=title,
                image_path=content.get("image_path", ""),
                alt_text=content.get("alt_text", ""),
                notes=notes
            )

        elif slide_type == "quote":
            self.add_quote_slide(
                quote_text=content.get("text", ""),
                attribution=content.get("attribution", ""),
                notes=notes
            )

        elif slide_type == "blank":
            self.add_blank_slide(notes=notes)

        elif slide_type == "styled_content":
            self.add_styled_content_slide(
                title=title,
                bullets=content.get("bullets", []),
                step_number=content.get("step_number"),
                callout_text=content.get("callout_text"),
                callout_type=content.get("callout_type", "tip"),
                notes=notes
            )

        elif slide_type == "table":
            self.add_table_slide(
                title=title,
                headers=content.get("headers", []),
                rows=content.get("rows", []),
                totals_row=content.get("totals_row"),
                notes=notes
            )

        elif slide_type == "process_flow":
            self.add_process_flow_slide(
                title=title,
                steps=content.get("steps", []),
                flow_style=content.get("flow_style", "horizontal"),
                notes=notes
            )

        elif slide_type == "comparison_matrix":
            self.add_comparison_matrix_slide(
                title=title,
                options=content.get("options", []),
                criteria=content.get("criteria", []),
                highlight_option=content.get("highlight_option"),
                notes=notes
            )

        elif slide_type == "timeline":
            self.add_timeline_slide(
                title=title,
                events=content.get("events", []),
                notes=notes
            )

        elif slide_type == "key_value":
            self.add_key_value_slide(
                title=title,
                pairs=content.get("pairs", []),
                notes=notes
            )

        elif slide_type == "callout":
            self.add_callout_slide(
                title=title,
                callouts=content.get("callouts", []),
                notes=notes
            )

        elif slide_type == "stats_grid":
            self.add_stats_grid_slide(
                title=title,
                stats=content.get("stats", []),
                layout=content.get("layout", "auto"),
                notes=notes
            )

        elif slide_type == "cards_grid":
            self.add_cards_grid_slide(
                title=title,
                cards=content.get("cards", []),
                columns=content.get("columns", 2),
                notes=notes
            )

        elif slide_type == "key_concept":
            self.add_key_concept_slide(
                title=title,
                concept_title=content.get("concept_title", ""),
                definition=content.get("definition", ""),
                details=content.get("details"),
                concept_style=content.get("concept_style", "boxed"),
                notes=notes
            )

        elif slide_type == "image_grid":
            self.add_image_grid_slide(slide_data)

        elif slide_type == "data_viz":
            self.add_data_viz_slide(slide_data)

        else:
            # Default to content slide
            self.add_content_slide(
                title=title,
                bullets=content.get("bullets", []),
                notes=notes
            )

    def _get_layout(self, slide_type: SlideType):
        """Get the slide layout for a given type with template-aware discovery."""
        # First check template-specific layout overrides
        layout_key = slide_type.value  # e.g., "title", "content"
        if layout_key in self._layout_overrides:
            idx = self._layout_overrides[layout_key]
        else:
            # Fall back to default indices
            idx = self.LAYOUT_INDICES.get(slide_type, 1)

        # Ensure index is valid for this presentation
        if idx >= len(self.prs.slide_layouts):
            logger.warning(f"Layout index {idx} not found for {slide_type.value}, using fallback")
            idx = 1  # Fallback to title and content

        return self.prs.slide_layouts[idx]

    def _add_notes(self, slide, notes: str) -> None:
        """Add speaker notes to a slide."""
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    def _fill_text_frame(self, text_frame, items: List[str]) -> None:
        """Fill a text frame with bullet items."""
        text_frame.clear()
        for i, item in enumerate(items):
            if i == 0:
                text_frame.paragraphs[0].text = item
            else:
                p = text_frame.add_paragraph()
                p.text = item
                p.level = 0

    def _fill_comparison_content(
        self,
        slide,
        left_title: str,
        left_content: List[str],
        right_title: str,
        right_content: List[str]
    ) -> None:
        """Fill comparison slide content."""
        # This handles various comparison layout structures
        placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

        # Try to find and fill the comparison placeholders
        # Standard comparison layout has indices: 0 (title), 1 (left text title),
        # 2 (left content), 3 (right text title), 4 (right content)

        if 1 in placeholders and left_title:
            placeholders[1].text = left_title
        if 2 in placeholders and left_content:
            self._fill_text_frame(placeholders[2].text_frame, left_content)
        if 3 in placeholders and right_title:
            placeholders[3].text = right_title
        if 4 in placeholders and right_content:
            self._fill_text_frame(placeholders[4].text_frame, right_content)

    def save(self, output_path: str) -> str:
        """
        Save the presentation to a file.

        Args:
            output_path: Path for the output .pptx file

        Returns:
            The path to the saved file
        """
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        self.prs.save(str(output_path))
        logger.info(f"Saved presentation: {output_path} ({self._slide_count} slides)")

        return str(output_path)

    @property
    def slide_count(self) -> int:
        """Get the current slide count."""
        return self._slide_count


def create_presentation_from_json(
    json_path: str,
    output_path: str,
    template_path: Optional[str] = None,
    template_name: Optional[str] = None
) -> str:
    """
    Create a presentation from a JSON file.

    Args:
        json_path: Path to JSON file with presentation content
        output_path: Path for output PPTX file
        template_path: Optional template PPTX file path
        template_name: Optional registered template name (e.g., "corporate", "dark_mode")

    Returns:
        Path to the created PPTX file
    """
    with open(json_path, 'r', encoding='utf-8') as f:
        content = json.load(f)

    generator = PPTXGenerator(template_path=template_path, template_name=template_name)
    generator.create_from_structure(content)

    return generator.save(output_path)


def create_presentation_from_markdown(
    markdown_text: str,
    output_path: str,
    template_path: Optional[str] = None
) -> str:
    """
    Create a presentation from markdown text.

    Uses a simple markdown structure:
    - # Title becomes title slide
    - ## Section becomes section header
    - ### Slide Title followed by bullet points becomes content slide

    Args:
        markdown_text: Markdown formatted text
        output_path: Path for output PPTX file
        template_path: Optional template PPTX file

    Returns:
        Path to the created PPTX file
    """
    generator = PPTXGenerator(template_path=template_path)

    lines = markdown_text.strip().split('\n')
    current_section = {"title": "", "slides": []}
    current_slide = None
    sections = []
    title_found = False

    for line in lines:
        line = line.strip()

        if not line:
            continue

        # Title (H1)
        if line.startswith('# ') and not title_found:
            title_found = True
            generator.add_title_slide(title=line[2:].strip())

        # Section header (H2)
        elif line.startswith('## '):
            if current_section["slides"]:
                sections.append(current_section)
            current_section = {"title": line[3:].strip(), "slides": []}
            generator.add_section_header(title=current_section["title"])

        # Slide title (H3)
        elif line.startswith('### '):
            if current_slide:
                generator.add_content_slide(
                    title=current_slide["title"],
                    bullets=current_slide["bullets"]
                )
            current_slide = {"title": line[4:].strip(), "bullets": []}

        # Bullet point
        elif line.startswith('- ') or line.startswith('* '):
            if current_slide is not None:
                current_slide["bullets"].append(line[2:].strip())

    # Add final slide
    if current_slide:
        generator.add_content_slide(
            title=current_slide["title"],
            bullets=current_slide["bullets"]
        )

    return generator.save(output_path)


# CLI interface
if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description="Generate PowerPoint presentations from structured content"
    )
    parser.add_argument(
        "--input", "-i",
        help="Input JSON file with presentation content"
    )
    parser.add_argument(
        "--output", "-o",
        help="Output PPTX file path"
    )
    parser.add_argument(
        "--template", "-t",
        help="Path to PPTX template file for theming"
    )
    parser.add_argument(
        "--theme",
        help="Registered template name (e.g., corporate, dark_mode, creative, minimal, educational)"
    )
    parser.add_argument(
        "--list-templates",
        action="store_true",
        help="List all available templates and exit"
    )

    args = parser.parse_args()

    # Handle --list-templates
    if args.list_templates:
        if TEMPLATE_LOADER_AVAILABLE:
            loader = TemplateLoader()
            print("\n=== Available Templates ===\n")
            for template in loader.list_templates():
                status = "OK" if template.path.exists() else "MISSING"
                print(f"  [{status}] {template.id}")
                print(f"         Name: {template.name}")
                print(f"         Category: {template.category}")
                print(f"         Description: {template.description}")
                print(f"         Primary color: {template.colors.primary}")
                print()
            print(f"Default template: {loader.get_default_template()}")
            print("\nUsage: python pptx_generator.py -i content.json -o output.pptx --theme corporate")
        else:
            print("Template loader not available. Use --template with a direct path instead.")
        sys.exit(0)

    # Validate required arguments for generation
    if not args.input or not args.output:
        parser.error("--input and --output are required for presentation generation")

    output = create_presentation_from_json(
        json_path=args.input,
        output_path=args.output,
        template_path=args.template,
        template_name=args.theme
    )

    print(f"Created presentation: {output}")
