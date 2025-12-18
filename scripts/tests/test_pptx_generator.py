#!/usr/bin/env python3
"""
Comprehensive unit tests for the PPTX Generator module.

Tests cover:
- Individual slide type methods
- Integration with create_from_structure
- Template and theme functionality
- Edge cases and error handling
- Output validation
"""

import pytest
import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

# Add pptx-generator to path
sys.path.insert(0, str(Path(__file__).parent.parent / "pptx-generator"))
from pptx_generator import PPTXGenerator, SlideType, PresentationMetadata, ThemeColors


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def validate_pptx_file(file_path: Path) -> bool:
    """Validate that a file is a valid PPTX file."""
    if not file_path.exists():
        return False
    if not file_path.suffix == '.pptx':
        return False
    try:
        prs = Presentation(str(file_path))
        return True
    except Exception:
        return False


def count_slides(file_path: Path) -> int:
    """Count slides in a PPTX file."""
    prs = Presentation(str(file_path))
    return len(prs.slides)


def get_slide_notes(file_path: Path, slide_index: int) -> str:
    """Get speaker notes from a specific slide."""
    prs = Presentation(str(file_path))
    slide = prs.slides[slide_index]
    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        return notes_frame.text
    return ""


# =============================================================================
# TEST CLASS: Slide Types
# =============================================================================

class TestSlideTypes:
    """Tests for individual slide type methods"""

    def test_add_title_slide(self, temp_output_dir):
        """Test creating a title slide"""
        generator = PPTXGenerator()
        generator.add_title_slide(
            title="Test Presentation",
            subtitle="A Test Subtitle",
            notes="Title slide notes"
        )

        output_path = temp_output_dir / "title_slide.pptx"
        generator.save(str(output_path))

        assert output_path.exists()
        assert validate_pptx_file(output_path)
        assert count_slides(output_path) == 1
        assert generator.slide_count == 1

    def test_add_section_header(self, temp_output_dir):
        """Test creating a section header slide"""
        generator = PPTXGenerator()
        generator.add_section_header(
            title="Section One",
            subtitle="Introduction to the topic",
            notes="Section header notes"
        )

        output_path = temp_output_dir / "section_header.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1
        assert generator.slide_count == 1

    def test_add_content_slide(self, temp_output_dir):
        """Test creating a standard content slide with bullets"""
        generator = PPTXGenerator()
        generator.add_content_slide(
            title="Key Points",
            bullets=["First point", "Second point", "Third point"],
            notes="Content slide notes"
        )

        output_path = temp_output_dir / "content_slide.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1
        assert generator.slide_count == 1

    def test_add_two_content_slide(self, temp_output_dir):
        """Test creating a two-column content slide"""
        generator = PPTXGenerator()
        generator.add_two_content_slide(
            title="Two Columns",
            left_content=["Left 1", "Left 2"],
            right_content=["Right 1", "Right 2"],
            notes="Two content notes"
        )

        output_path = temp_output_dir / "two_content.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1

    def test_add_comparison_slide(self, temp_output_dir):
        """Test creating a comparison slide"""
        generator = PPTXGenerator()
        generator.add_comparison_slide(
            title="Comparison",
            left_title="Pros",
            left_content=["Advantage 1", "Advantage 2"],
            right_title="Cons",
            right_content=["Disadvantage 1", "Disadvantage 2"],
            notes="Comparison notes"
        )

        output_path = temp_output_dir / "comparison.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1

    def test_add_image_slide(self, temp_output_dir):
        """Test creating an image slide (without actual image file)"""
        generator = PPTXGenerator()
        generator.add_image_slide(
            title="Image Display",
            image_path="/path/to/nonexistent/image.png",
            alt_text="Sample image description",
            notes="Image slide notes"
        )

        output_path = temp_output_dir / "image_slide.pptx"
        generator.save(str(output_path))

        # Note: When image doesn't exist, slide may not be created
        # This tests graceful handling of missing images
        assert validate_pptx_file(output_path)

    def test_add_quote_slide(self, temp_output_dir):
        """Test creating a quote slide"""
        generator = PPTXGenerator()
        generator.add_quote_slide(
            quote_text="The best way to predict the future is to create it.",
            attribution="Peter Drucker",
            notes="Quote slide notes"
        )

        output_path = temp_output_dir / "quote.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1

    def test_add_table_slide(self, temp_output_dir):
        """Test creating a table slide"""
        generator = PPTXGenerator()
        generator.add_table_slide(
            title="Data Table",
            headers=["Feature", "Basic", "Pro"],
            rows=[
                ["Users", "5", "25"],
                ["Storage", "10GB", "100GB"]
            ],
            totals_row=["Total", "$10", "$50"],
            notes="Table slide notes"
        )

        output_path = temp_output_dir / "table.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1

    def test_add_process_flow_slide(self, temp_output_dir):
        """Test creating a process flow slide"""
        generator = PPTXGenerator()
        generator.add_process_flow_slide(
            title="Process Steps",
            steps=[
                {"label": "Planning", "description": "Define requirements"},
                {"label": "Development", "description": "Build solution"},
                {"label": "Testing", "description": "Quality assurance"}
            ],
            flow_style="horizontal",
            notes="Process flow notes"
        )

        output_path = temp_output_dir / "process_flow.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1

    def test_add_timeline_slide(self, temp_output_dir):
        """Test creating a timeline slide"""
        generator = PPTXGenerator()
        generator.add_timeline_slide(
            title="Project Timeline",
            events=[
                {"date": "Jan 2025", "label": "Kickoff", "description": "Planning"},
                {"date": "Mar 2025", "label": "Alpha", "description": "Testing"}
            ],
            notes="Timeline notes"
        )

        output_path = temp_output_dir / "timeline.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1

    def test_add_key_value_slide(self, temp_output_dir):
        """Test creating a key-value pairs slide"""
        generator = PPTXGenerator()
        generator.add_key_value_slide(
            title="Configuration",
            pairs=[
                {"key": "Environment", "value": "Production"},
                {"key": "Region", "value": "US-East-1"}
            ],
            notes="Key-value notes"
        )

        output_path = temp_output_dir / "key_value.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1

    def test_add_callout_slide(self, temp_output_dir):
        """Test creating a callout boxes slide"""
        generator = PPTXGenerator()
        generator.add_callout_slide(
            title="Important Information",
            callouts=[
                {"callout_type": "info", "heading": "Info", "text": "Information here"},
                {"callout_type": "warning", "heading": "Warning", "text": "Be careful"}
            ],
            notes="Callout notes"
        )

        output_path = temp_output_dir / "callout.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1

    def test_add_stats_grid_slide(self, temp_output_dir):
        """Test creating a statistics grid slide"""
        generator = PPTXGenerator()
        generator.add_stats_grid_slide(
            title="Key Metrics",
            stats=[
                {"value": "95%", "label": "Satisfaction", "trend": "up", "highlight": True},
                {"value": "$1.2M", "label": "Revenue", "trend": "up", "highlight": False}
            ],
            layout="auto",
            notes="Stats grid notes"
        )

        output_path = temp_output_dir / "stats_grid.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1

    def test_add_cards_grid_slide(self, temp_output_dir):
        """Test creating a cards grid slide"""
        generator = PPTXGenerator()
        generator.add_cards_grid_slide(
            title="Feature Categories",
            cards=[
                {"title": "Performance", "items": ["Fast", "Optimized"], "color": "primary"},
                {"title": "Security", "items": ["Encrypted", "Secure"], "color": "secondary"}
            ],
            columns=2,
            notes="Cards grid notes"
        )

        output_path = temp_output_dir / "cards_grid.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1

    def test_add_data_viz_slide(self, temp_output_dir):
        """Test creating a data visualization (chart) slide"""
        generator = PPTXGenerator()
        slide_data = {
            "title": "Sales Data",
            "content": {
                "chart_type": "column_clustered",
                "chart_data": {
                    "categories": ["Q1", "Q2", "Q3"],
                    "series": [
                        {"name": "2024", "values": [10, 20, 15]}
                    ]
                },
                "chart_title": "Quarterly Sales",
                "show_legend": True,
                "show_data_labels": False
            },
            "notes": "Chart notes"
        }
        generator.add_data_viz_slide(slide_data)

        output_path = temp_output_dir / "data_viz.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1

    def test_add_divider_slide(self, temp_output_dir):
        """Test creating a divider slide"""
        generator = PPTXGenerator()
        slide_data = {
            "title": "Next Section",
            "content": {
                "section_number": "02",
                "subtitle": "Introduction",
                "divider_style": "bold",
                "accent_color": "primary"
            },
            "notes": "Divider notes"
        }
        generator.add_divider_slide(slide_data)

        output_path = temp_output_dir / "divider.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1

    def test_add_agenda_slide(self, temp_output_dir):
        """Test creating an agenda slide"""
        generator = PPTXGenerator()
        slide_data = {
            "title": "Today's Agenda",
            "content": {
                "agenda_items": [
                    {"section": "Introduction", "description": "10 min", "status": "complete"},
                    {"section": "Main Topic", "description": "30 min", "status": "current"}
                ],
                "agenda_style": "numbered",
                "show_progress": True
            },
            "notes": "Agenda notes"
        }
        generator.add_agenda_slide(slide_data)

        output_path = temp_output_dir / "agenda.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1

    def test_add_image_grid_slide(self, temp_output_dir):
        """Test creating an image grid slide"""
        generator = PPTXGenerator()
        slide_data = {
            "title": "Image Gallery",
            "content": {
                "images": [
                    {"path": "/path/to/img1.png", "caption": "Image 1", "alt_text": "First"},
                    {"path": "/path/to/img2.png", "caption": "Image 2", "alt_text": "Second"}
                ],
                "grid_layout": "auto",
                "show_captions": True,
                "gap_size": "medium"
            },
            "notes": "Image grid notes"
        }
        generator.add_image_grid_slide(slide_data)

        output_path = temp_output_dir / "image_grid.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1

    def test_add_blank_slide(self, temp_output_dir):
        """Test creating a blank slide"""
        generator = PPTXGenerator()
        generator.add_blank_slide(notes="Blank slide notes")

        output_path = temp_output_dir / "blank.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1


# =============================================================================
# TEST CLASS: Integration Tests
# =============================================================================

class TestCreateFromStructure:
    """Tests for create_from_structure method"""

    def test_create_from_structure_minimal(self, temp_output_dir, minimal_presentation):
        """Test creating presentation from minimal 5-slide fixture"""
        generator = PPTXGenerator()
        generator.create_from_structure(minimal_presentation)

        output_path = temp_output_dir / "minimal.pptx"
        generator.save(str(output_path))

        assert validate_pptx_file(output_path)
        assert count_slides(output_path) == 5
        assert generator.slide_count == 5

    def test_create_from_structure_all_types(self, temp_output_dir, all_slide_types_presentation):
        """Test creating presentation with all slide types"""
        generator = PPTXGenerator()
        generator.create_from_structure(all_slide_types_presentation)

        output_path = temp_output_dir / "all_types.pptx"
        generator.save(str(output_path))

        assert validate_pptx_file(output_path)
        # Should have multiple slides covering all types
        assert generator.slide_count > 15

    def test_create_from_structure_complex(self, temp_output_dir, complex_presentation):
        """Test creating large 50-slide presentation"""
        generator = PPTXGenerator()
        generator.create_from_structure(complex_presentation)

        output_path = temp_output_dir / "complex.pptx"
        generator.save(str(output_path))

        assert validate_pptx_file(output_path)
        assert generator.slide_count == 50

    def test_metadata_extraction(self, temp_output_dir, minimal_presentation):
        """Test that metadata is properly set from structure"""
        generator = PPTXGenerator()
        generator.create_from_structure(minimal_presentation)

        assert generator.metadata.title == "Minimal Test Presentation"
        assert generator.metadata.subtitle == "Five Slide Sample for Testing"
        assert generator.metadata.author == "Slideforge Test Suite"

    def test_speaker_notes_included(self, temp_output_dir, minimal_presentation):
        """Test that speaker notes are included in slides"""
        generator = PPTXGenerator()
        generator.create_from_structure(minimal_presentation)

        output_path = temp_output_dir / "notes_test.pptx"
        generator.save(str(output_path))

        # Check first slide has notes
        notes = get_slide_notes(output_path, 0)
        assert "title slide" in notes.lower()


# =============================================================================
# TEST CLASS: Templates and Themes
# =============================================================================

class TestTemplates:
    """Tests for template and theme functionality"""

    def test_default_theme_colors(self):
        """Test that default theme colors are applied"""
        generator = PPTXGenerator()

        assert generator.theme_colors.primary == "#2c5aa0"
        assert generator.theme_colors.secondary == "#28a745"
        assert generator.theme_colors.accent == "#ffc107"

    def test_custom_theme_colors(self):
        """Test custom theme color application"""
        custom_colors = ThemeColors(
            primary="#ff0000",
            secondary="#00ff00",
            accent="#0000ff"
        )
        generator = PPTXGenerator(theme_colors=custom_colors)

        assert generator.theme_colors.primary == "#ff0000"
        assert generator.theme_colors.secondary == "#00ff00"
        assert generator.theme_colors.accent == "#0000ff"

    def test_template_path_nonexistent(self):
        """Test graceful handling of non-existent template path"""
        # Should not raise exception, just use default
        generator = PPTXGenerator(template_path="/nonexistent/template.pptx")
        assert generator.prs is not None

    def test_metadata_setting(self):
        """Test setting presentation metadata"""
        generator = PPTXGenerator()
        metadata = PresentationMetadata(
            title="Test Title",
            author="Test Author",
            subject="Test Subject"
        )
        generator.set_metadata(metadata)

        assert generator.metadata.title == "Test Title"
        assert generator.metadata.author == "Test Author"


# =============================================================================
# TEST CLASS: Edge Cases
# =============================================================================

class TestEdgeCases:
    """Tests for edge cases and error handling"""

    def test_empty_content_list(self, temp_output_dir):
        """Test handling of empty bullet lists"""
        generator = PPTXGenerator()
        generator.add_content_slide(
            title="Empty Content",
            bullets=[],
            notes=""
        )

        output_path = temp_output_dir / "empty_bullets.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1

    def test_long_text_in_bullets(self, temp_output_dir):
        """Test handling of very long bullet text"""
        generator = PPTXGenerator()
        long_text = "A" * 500  # 500 character string
        generator.add_content_slide(
            title="Long Text",
            bullets=[long_text, "Normal text"],
            notes=""
        )

        output_path = temp_output_dir / "long_text.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1

    def test_special_characters_in_text(self, temp_output_dir):
        """Test handling of special characters"""
        generator = PPTXGenerator()
        generator.add_content_slide(
            title="Special Chars: <>&\"'",
            bullets=[
                "Ampersand: &",
                "Less than: <",
                "Greater than: >",
                "Quotes: \"'",
                "Unicode: \u00e9\u00f1\u00fc"
            ],
            notes="Notes with special chars: <>&\"'"
        )

        output_path = temp_output_dir / "special_chars.pptx"
        generator.save(str(output_path))

        assert validate_pptx_file(output_path)

    def test_missing_optional_fields(self, temp_output_dir):
        """Test handling of missing optional content fields"""
        minimal_structure = {
            "metadata": {
                "title": "Minimal"
            },
            "sections": [
                {
                    "slides": [
                        {
                            "type": "title",
                            "title": "Test"
                        }
                    ]
                }
            ]
        }

        generator = PPTXGenerator()
        generator.create_from_structure(minimal_structure)

        output_path = temp_output_dir / "minimal_fields.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1

    def test_max_bullets_enforcement(self, temp_output_dir):
        """Test handling of more than 6 bullets (6x6 rule)"""
        generator = PPTXGenerator()
        # Create 10 bullets (exceeds recommended 6)
        many_bullets = [f"Bullet {i+1}" for i in range(10)]
        generator.add_content_slide(
            title="Too Many Bullets",
            bullets=many_bullets,
            notes=""
        )

        output_path = temp_output_dir / "many_bullets.pptx"
        generator.save(str(output_path))

        # Should still create slide (no hard limit, just recommendation)
        assert count_slides(output_path) == 1

    def test_empty_presentation(self, temp_output_dir):
        """Test saving presentation with no slides"""
        generator = PPTXGenerator()

        output_path = temp_output_dir / "empty.pptx"
        generator.save(str(output_path))

        assert validate_pptx_file(output_path)
        assert count_slides(output_path) == 0

    def test_unicode_in_title_and_content(self, temp_output_dir):
        """Test handling of Unicode characters in various languages"""
        generator = PPTXGenerator()
        generator.add_content_slide(
            title="Multilingual: 中文 日本語 العربية",
            bullets=[
                "Chinese: 你好",
                "Japanese: こんにちは",
                "Arabic: مرحبا",
                "Greek: Γειά σου",
                "Russian: Привет"
            ],
            notes="Unicode notes"
        )

        output_path = temp_output_dir / "unicode.pptx"
        generator.save(str(output_path))

        assert validate_pptx_file(output_path)

    def test_table_with_empty_cells(self, temp_output_dir):
        """Test table slide with empty cells"""
        generator = PPTXGenerator()
        generator.add_table_slide(
            title="Table with Gaps",
            headers=["Col1", "Col2", "Col3"],
            rows=[
                ["Value", "", "Value"],
                ["", "Value", ""],
                ["Value", "Value", "Value"]
            ],
            notes=""
        )

        output_path = temp_output_dir / "table_gaps.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1

    def test_none_values_in_content(self, temp_output_dir):
        """Test handling of None values in optional fields"""
        generator = PPTXGenerator()
        generator.add_content_slide(
            title="Test",
            bullets=["Item 1", "Item 2"],
            notes=None  # None instead of empty string
        )

        output_path = temp_output_dir / "none_values.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1


# =============================================================================
# TEST CLASS: Output Validation
# =============================================================================

class TestOutputValidation:
    """Tests for output file validation"""

    def test_pptx_file_valid(self, temp_output_dir, minimal_presentation):
        """Verify output is a valid PPTX file"""
        generator = PPTXGenerator()
        generator.create_from_structure(minimal_presentation)

        output_path = temp_output_dir / "valid_test.pptx"
        generator.save(str(output_path))

        # Should be readable by python-pptx
        assert validate_pptx_file(output_path)

        # Should be able to load and re-save
        prs = Presentation(str(output_path))
        resave_path = temp_output_dir / "resaved.pptx"
        prs.save(str(resave_path))
        assert resave_path.exists()

    def test_slide_count_accuracy(self, temp_output_dir):
        """Verify slide count property matches actual slides"""
        generator = PPTXGenerator()

        # Add various slide types
        generator.add_title_slide("Title", "Subtitle")
        generator.add_content_slide("Content", ["Bullet 1"])
        generator.add_section_header("Section", "Description")

        assert generator.slide_count == 3

        output_path = temp_output_dir / "count_test.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 3

    def test_speaker_notes_persistence(self, temp_output_dir):
        """Verify speaker notes are saved correctly"""
        generator = PPTXGenerator()
        test_notes = "These are important speaker notes for testing."

        generator.add_content_slide(
            title="Test Slide",
            bullets=["Point 1", "Point 2"],
            notes=test_notes
        )

        output_path = temp_output_dir / "notes_persistence.pptx"
        generator.save(str(output_path))

        notes = get_slide_notes(output_path, 0)
        assert test_notes in notes

    def test_output_directory_creation(self, tmp_path):
        """Test that output directory is created if it doesn't exist"""
        generator = PPTXGenerator()
        generator.add_title_slide("Test", "Test")

        # Create path in non-existent directory
        deep_path = tmp_path / "level1" / "level2" / "level3" / "output.pptx"
        generator.save(str(deep_path))

        assert deep_path.exists()
        assert validate_pptx_file(deep_path)

    def test_file_extension_enforcement(self, temp_output_dir):
        """Test that .pptx extension is handled correctly"""
        generator = PPTXGenerator()
        generator.add_title_slide("Test", "Test")

        # Try saving without extension
        output_path = temp_output_dir / "no_extension"
        saved_path = generator.save(str(output_path))

        # File should still be created (python-pptx handles this)
        assert Path(saved_path).exists()

    def test_consecutive_presentations(self, temp_output_dir):
        """Test creating multiple presentations in sequence"""
        for i in range(3):
            generator = PPTXGenerator()
            generator.add_title_slide(f"Presentation {i+1}", f"Number {i+1}")

            output_path = temp_output_dir / f"pres_{i+1}.pptx"
            generator.save(str(output_path))

            assert validate_pptx_file(output_path)
            assert count_slides(output_path) == 1


# =============================================================================
# TEST CLASS: Comparison Matrix
# =============================================================================

class TestComparisonMatrix:
    """Tests for comparison matrix slide type"""

    def test_add_comparison_matrix_slide(self, temp_output_dir):
        """Test creating a comparison matrix slide"""
        generator = PPTXGenerator()
        generator.add_comparison_matrix_slide(
            title="Product Comparison",
            options=["Option A", "Option B", "Option C"],
            criteria=[
                {"name": "Price", "values": ["$10", "$20", "$30"]},
                {"name": "Features", "values": ["Basic", "Standard", "Premium"]}
            ],
            highlight_option="Option B",
            notes="Comparison matrix notes"
        )

        output_path = temp_output_dir / "comparison_matrix.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1


# =============================================================================
# TEST CLASS: Key Concept
# =============================================================================

class TestKeyConcept:
    """Tests for key concept slide type"""

    def test_add_key_concept_slide(self, temp_output_dir):
        """Test creating a key concept slide"""
        generator = PPTXGenerator()
        generator.add_key_concept_slide(
            title="Important Concept",
            concept_title="Microservices",
            definition="An architectural style for building applications",
            details=["Independently deployable", "Loosely coupled", "Technology agnostic"],
            concept_style="boxed",
            notes="Key concept notes"
        )

        output_path = temp_output_dir / "key_concept.pptx"
        generator.save(str(output_path))

        assert count_slides(output_path) == 1


# =============================================================================
# RUN TESTS
# =============================================================================

if __name__ == "__main__":
    pytest.main([__file__, "-v", "--tb=short"])
